# -*- coding: utf-8 -*-
"""Generate a PPT showing the 9 machine-verified 7x7 StandardDraw positions,
one slide per position, each with its board image (PNG) and verification facts."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ---------------- palette ----------------
NAVY    = RGBColor(0x1B, 0x2A, 0x4A)
NAVY2   = RGBColor(0x2C, 0x3E, 0x63)
BLUE    = RGBColor(0x3B, 0x59, 0x98)
BLUE_L  = RGBColor(0xE7, 0xEC, 0xF6)
GOLD    = RGBColor(0xE8, 0xA3, 0x3D)
GOLD_D  = RGBColor(0xC9, 0x8A, 0x2D)
GOLD_L  = RGBColor(0xFB, 0xEF, 0xDA)
GREEN   = RGBColor(0x2E, 0x7D, 0x4F)
GREEN_L = RGBColor(0xE4, 0xF1, 0xE9)
RED     = RGBColor(0xB3, 0x40, 0x2F)
RED_L   = RGBColor(0xF8, 0xE8, 0xE4)
INK     = RGBColor(0x2B, 0x2B, 0x2B)
GRAY    = RGBColor(0x6B, 0x74, 0x80)
LIGHT   = RGBColor(0xF1, 0xF3, 0xF7)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
FONT    = '微软雅黑'

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
BOARD_DIR = os.path.join(ROOT, 'docs', 'boards')

# 9 positions: (lean file, seed label, board png, white cert nodes, black cert nodes)
POSITIONS = [
    ('Draw7x7.lean',   'seed 1', 'Draw7x7.png',   400, 158),
    ('Draw7x7s2.lean', 'seed 2', 'Draw7x7s2.png', None, None),
    ('Draw7x7s3.lean', 'seed 3', 'Draw7x7s3.png', None, None),
    ('Draw7x7s4.lean', 'seed 4', 'Draw7x7s4.png', None, None),
    ('Draw7x7s5.lean', 'seed 5', 'Draw7x7s5.png', None, None),
    ('Draw7x7s6.lean', 'seed 6', 'Draw7x7s6.png', None, None),
    ('Draw7x7s7.lean', 'seed 7', 'Draw7x7s7.png', None, None),
    ('Draw7x7s8.lean', 'seed 8', 'Draw7x7s8.png', None, None),
    ('Draw7x7s9.lean', 'seed 9', 'Draw7x7s9.png', None, None),
]
TOTAL = 1 + 1 + len(POSITIONS) + 1  # cover + overview + 9 positions + summary

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def new_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def _set_ea(run, name=FONT):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', name)


def text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = para.get('align', align)
        if para.get('space_before') is not None:
            p.space_before = Pt(para['space_before'])
        if para.get('space_after') is not None:
            p.space_after = Pt(para['space_after'])
        if para.get('line_spacing') is not None:
            p.line_spacing = para['line_spacing']
        for r in para['runs']:
            run = p.add_run()
            run.text = r[0]
            f = run.font
            f.name = FONT
            f.size = Pt(r[1])
            f.bold = r[2] if len(r) > 2 else False
            f.color.rgb = r[3] if len(r) > 3 else INK
            _set_ea(run)
    return tb


def box(slide, x, y, w, h, fill=None, line=None, lw=0.75,
        shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    if radius is not None:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def text_in(sp, paras, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = para.get('align', align)
        if para.get('space_after') is not None:
            p.space_after = Pt(para['space_after'])
        for r in para['runs']:
            run = p.add_run()
            run.text = r[0]
            f = run.font
            f.name = FONT
            f.size = Pt(r[1])
            f.bold = r[2] if len(r) > 2 else False
            f.color.rgb = r[3] if len(r) > 3 else INK
            _set_ea(run)
    return sp


def header(slide, title, kicker=None):
    box(slide, 0, 0, 13.333, 1.0, fill=NAVY)
    box(slide, 0.62, 0.33, 0.11, 0.36, fill=GOLD)
    text(slide, 0.9, 0.16, 10.5, 0.68,
         [{'runs': [(title, 28, True, WHITE)]}], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        text(slide, 8.7, 0.3, 4.0, 0.4,
             [{'runs': [(kicker, 12, False, RGBColor(0xB9, 0xC4, 0xDD))],
               'align': PP_ALIGN.RIGHT}], anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, idx):
    text(slide, 0.62, 7.12, 6, 0.3,
         [{'runs': [('Gomoku Formalization · 9 个机器验证的和棋局面', 9, False, GRAY)]}])
    text(slide, 11.9, 7.12, 0.85, 0.3,
         [{'runs': [(f'{idx} / {TOTAL}', 9, False, GRAY)],
           'align': PP_ALIGN.RIGHT}])


def check_chip(slide, x, y, w, ok, txt, size=13.5):
    fill = GREEN_L if ok else RED_L
    color = GREEN if ok else RED
    sp = box(slide, x, y, w, 0.5, fill=fill,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    text_in(sp, [{'runs': [(('✓ ' if ok else '✗ ') + txt, size, True, color)]}])
    return sp


# ================================================================
# Slide 1 · Cover
# ================================================================
s = new_slide(NAVY)
box(s, 0, 7.28, 13.333, 0.22, fill=GOLD)
text(s, 0.9, 1.1, 11.5, 0.45,
     [{'runs': [('GOMOKU FORMALIZATION · VERIFIED DRAW POSITIONS', 14, True, GOLD)]}])
text(s, 0.9, 1.75, 11.5, 1.3,
     [{'runs': [('9 个机器验证的 7×7 和棋局面', 44, True, WHITE)]}])
text(s, 0.9, 3.2, 11.0, 0.6,
     [{'runs': [('每个局面: 21 黑 + 21 白 · 轮到黑方 · 7 个空格 · 双方都无法成五',
                 18, False, RGBColor(0xCF, 0xD8, 0xE8))]}])
box(s, 0.95, 4.05, 2.6, 0.05, fill=GOLD)
text(s, 0.9, 4.3, 11.5, 0.55,
     [{'runs': [('不可信搜索器找证书  →  Lean 检查器逐节点重算  →  StandardDraw 定理',
                 17, False, RGBColor(0xF2, 0xCE, 0x94))]}])
text(s, 0.9, 6.3, 11.5, 0.4,
     [{'runs': [('C++ DefenseSearcher · checkDefenseCertificateAt · '
                 'white/black_defense_certificate_sound · standardDraw_of_mutualDefense',
                 12, False, RGBColor(0x8E, 0x9B, 0xB8))]}])

# ================================================================
# Slide 2 · How a draw is verified
# ================================================================
s = new_slide()
header(s, '和棋是怎么验证的', 'Trust Chain')
steps = [
    ('① 构造局面', '每个长度 5 窗口都含黑白两子\n→ 任何一方永远无法成五', NAVY2),
    ('② C++ 找证书', '白防黑 + 黑防白 两张防守证书\n(搜索器不可信, 只出候选)', BLUE),
    ('③ Lean 重算', 'checkDefenseCertificateAt\n逐节点重算局面/轮次/终局/应手', GOLD_D),
    ('④ 组合出定理', '两张证书 + soundness\n⇒ StandardDraw 局面', GREEN),
]
x = 0.62
for i, (t1, t2, c) in enumerate(steps):
    sp = box(s, x, 1.7, 2.85, 1.55, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.1)
    text_in(sp, [{'runs': [(t1, 17, True, WHITE)]},
                 {'runs': [(t2, 11.5, False, WHITE)]}])
    if i < 3:
        from pptx.enum.shapes import MSO_CONNECTOR
        conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x + 2.85), Inches(2.47),
                                      Inches(x + 3.28), Inches(2.47))
        conn.line.color.rgb = GRAY
        conn.line.width = Pt(2.0)
    x += 3.28
text(s, 0.62, 3.75, 12.1, 0.55,
     [{'runs': [('每个局面独立验证 4 项断言 (全部 native_decide, 不依赖 C++ 结果):',
                 16, True, NAVY)]}])
checks = [
    ('白防黑证书 checkDefenseCertificateAt = true', '9 / 9 通过'),
    ('黑防白证书 checkDefenseCertificateAt = true', '9 / 9 通过'),
    ('局面形状: 21黑 + 21白 + 7空 + 轮到黑方', '9 / 9 通过'),
    ('60 个长度 5 窗口全部双色', '9 / 9 通过'),
    ('StandardDraw <position> 由两张证书组合', '9 / 9 通过'),
]
y = 4.4
for i, (t, r) in enumerate(checks):
    box(s, 0.62, y, 8.6, 0.42, fill=LIGHT,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    text(s, 0.95, y, 8.2, 0.42, [{'runs': [(t, 14, False, INK)]}],
         anchor=MSO_ANCHOR.MIDDLE)
    chip = box(s, 9.5, y, 3.0, 0.42, fill=GREEN_L,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    text_in(chip, [{'runs': [(r, 13.5, True, GREEN)]}])
    y += 0.5
note = box(s, 0.62, 6.5, 12.09, 0.5, fill=GOLD_L,
           shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.2)
text_in(note, [{'runs': [('搜索器可以犯错, 检查器不放错; 通过检查的证书才是定理。',
                          14, True, GOLD_D)]}])
footer(s, 2)

# ================================================================
# Slides 3..11 · one per position
# ================================================================
for idx, (lean, seed, png, wnodes, bnodes) in enumerate(POSITIONS):
    page = 3 + idx
    s = new_slide()
    header(s, f'和棋局面 {idx + 1} / 9', f'{lean} · {seed}')

    # Board image on the left.
    img_path = os.path.join(BOARD_DIR, png)
    if os.path.exists(img_path):
        s.shapes.add_picture(img_path, Inches(0.75), Inches(1.35),
                             height=Inches(5.35))

    # Facts on the right.
    fx = 6.45
    facts = [
        ('局面构成', '21 黑 + 21 白 · 7 空格 · 轮到黑方'),
        ('结构性保证', '每个长度 5 窗口同时含黑白两子\n任何一方都无法成五'),
        ('白防黑证书', 'WhiteCanPreventBlackWin  √ 检查器通过'),
        ('黑防白证书', 'BlackCanPreventWhiteWin  √ 检查器通过'),
        ('和棋定理', 'StandardDraw <root>  √ 组合成立'),
    ]
    y = 1.5
    for i, (t1, t2) in enumerate(facts):
        box(s, fx, y, 6.1, 1.02, fill=LIGHT,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
        text(s, fx + 0.3, y + 0.12, 2.3, 0.4,
             [{'runs': [(t1, 15.5, True, BLUE)]}])
        text(s, fx + 0.3, y + 0.48, 5.6, 0.5,
             [{'runs': [(t2, 13, False, INK)], 'line_spacing': 1.1}])
        y += 1.13

    footer(s, page)

# ================================================================
# Slide 12 · Summary table
# ================================================================
s = new_slide()
header(s, '全部 9 个局面: 汇总', 'Summary')
col_x = [0.62, 3.1, 5.2, 7.0, 8.9, 10.9]
col_w = [2.48, 2.1, 1.8, 1.9, 2.0, 1.6]
heads = ['Lean 定理文件', '种子', '白防证书', '黑防证书', 'StandardDraw', '窗口双色']
for i, h in enumerate(heads):
    box(s, col_x[i], 1.5, col_w[i], 0.55, fill=NAVY)
    text(s, col_x[i], 1.5, col_w[i], 0.55,
         [{'runs': [(h, 14, True, WHITE)], 'align': PP_ALIGN.CENTER}],
         anchor=MSO_ANCHOR.MIDDLE)
y = 2.15
for idx, (lean, seed, png, wnodes, bnodes) in enumerate(POSITIONS):
    fill = LIGHT if idx % 2 == 0 else WHITE
    for i, val in enumerate([lean, seed, '✓', '✓', '✓ 定理', '✓']):
        box(s, col_x[i], y, col_w[i], 0.52, fill=fill)
        color = INK
        if val == '✓':
            color = GREEN
        elif val == '✓ 定理':
            color = GREEN
        text(s, col_x[i], y, col_w[i], 0.52,
             [{'runs': [(val, 13, (val != lean), color)],
               'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.52
band = box(s, 0.62, 6.35, 12.09, 0.6, fill=GREEN_L,
           shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
text_in(band, [{'runs': [('9 / 9 全部通过: 每张证书经 Lean 检查器重算, 每个局面都证明了 '
                          'StandardDraw。', 15, True, GREEN)]}])
footer(s, 12)

# ---------------- speaker notes ----------------
NOTES = [
    '这份 PPT 展示最近完成的工作: 9 个机器验证的 7×7 和棋局面。核心纪律不变: '
    'C++ 搜索器只负责找候选证书, 是否成立完全由 Lean 检查器决定。',
    '每个局面的验证分四步: 先构造一个每个长度5窗口都同时含黑白两子的局面, '
    '然后 C++ 分别搜索白防黑和黑防白两张防守证书, Lean 检查器逐节点重算, '
    '最后用 soundness 定理组合出 StandardDraw。九个局面、五项断言全部通过。',
    '局面一: Draw7x7, 最初的种子。21黑21白, 7个空格, 轮到黑方。',
    '局面二: Draw7x7s2, 随机化移除顺序得到的第二个种子。',
    '局面三: Draw7x7s3。',
    '局面四: Draw7x7s4。',
    '局面五: Draw7x7s5。',
    '局面六: Draw7x7s6。',
    '局面七: Draw7x7s7。',
    '局面八: Draw7x7s8。',
    '局面九: Draw7x7s9。',
    '汇总: 九个局面的两张证书全部通过 Lean 检查, 每个局面都得到了 StandardDraw 定理。'
    '需要说明: 这些是特定的中盘局面, 空棋盘的和棋定理仍然是开放目标。',
]
for slide, note in zip(prs.slides, NOTES):
    slide.notes_slide.notes_text_frame.text = note

prs.core_properties.title = '9 个机器验证的 7×7 和棋局面'
prs.core_properties.author = 'Gomoku Formalization'

OUT = os.path.join(ROOT, '9个和棋局面.pptx')
try:
    prs.save(OUT)
    print('saved:', OUT, 'slides =', len(prs.slides._sldIdLst))
except PermissionError:
    OUT = os.path.join(ROOT, '9个和棋局面_新版.pptx')
    prs.save(OUT)
    print('original file locked; saved as:', OUT, 'slides =', len(prs.slides._sldIdLst))

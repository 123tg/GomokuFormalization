# -*- coding: utf-8 -*-
"""Generate an INTERLEAVED report PPT for the Gomoku project.

Unlike the chapter-style deck (make_pptx.py, 28 slides), this version threads
the verified draw positions INTO the technical narrative:
  semantics (force-win AND/OR vs defense) -> two-mode searcher -> certificate
  protocols -> trust boundary -> live examples (9 draw positions, 3 per slide)
  -> verification chain -> summary -> gaps.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------------- palette ----------------
NAVY    = RGBColor(0x1B, 0x2A, 0x4A)
NAVY2   = RGBColor(0x2C, 0x3E, 0x63)
BLUE    = RGBColor(0x3B, 0x59, 0x98)
BLUE_L  = RGBColor(0xE7, 0xEC, 0xF6)
GOLD    = RGBColor(0xE8, 0xA3, 0x3D)
GOLD_D  = RGBColor(0xC9, 0x8A, 0x2D)
GOLD_L  = RGBColor(0xFB, 0xEF, 0xDA)
GREEN   = RGBColor(0x2E, 0x7D, 0x4F)
GREEN_L = RGBColor(0xE4, 0xF1, 0xE9)
RED     = RGBColor(0xB3, 0x40, 0x2F)
RED_L   = RGBColor(0xF8, 0xE8, 0xE4)
INK     = RGBColor(0x2B, 0x2B, 0x2B)
GRAY    = RGBColor(0x6B, 0x74, 0x80)
LIGHT   = RGBColor(0xF1, 0xF3, 0xF7)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
FONT    = '微软雅黑'

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
BOARD_DIR = os.path.join(REPO_ROOT, 'docs', 'boards')

DRAW_POSITIONS = [
    ('Draw7x7.lean',   'seed 1', 'Draw7x7.png'),
    ('Draw7x7s2.lean', 'seed 2', 'Draw7x7s2.png'),
    ('Draw7x7s3.lean', 'seed 3', 'Draw7x7s3.png'),
    ('Draw7x7s4.lean', 'seed 4', 'Draw7x7s4.png'),
    ('Draw7x7s5.lean', 'seed 5', 'Draw7x7s5.png'),
    ('Draw7x7s6.lean', 'seed 6', 'Draw7x7s6.png'),
    ('Draw7x7s7.lean', 'seed 7', 'Draw7x7s7.png'),
    ('Draw7x7s8.lean', 'seed 8', 'Draw7x7s8.png'),
    ('Draw7x7s9.lean', 'seed 9', 'Draw7x7s9.png'),
]
TOTAL = 20

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def new_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def _set_ea(run, name=FONT):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', name)


def text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = para.get('align', align)
        if para.get('space_before') is not None:
            p.space_before = Pt(para['space_before'])
        if para.get('space_after') is not None:
            p.space_after = Pt(para['space_after'])
        if para.get('line_spacing') is not None:
            p.line_spacing = para['line_spacing']
        for r in para['runs']:
            run = p.add_run()
            run.text = r[0]
            f = run.font
            f.name = FONT
            f.size = Pt(r[1])
            f.bold = r[2] if len(r) > 2 else False
            f.color.rgb = r[3] if len(r) > 3 else INK
            _set_ea(run)
    return tb


def box(slide, x, y, w, h, fill=None, line=None, lw=0.75,
        shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    if radius is not None:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def text_in(sp, paras, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = para.get('align', align)
        if para.get('space_after') is not None:
            p.space_after = Pt(para['space_after'])
        for r in para['runs']:
            run = p.add_run()
            run.text = r[0]
            f = run.font
            f.name = FONT
            f.size = Pt(r[1])
            f.bold = r[2] if len(r) > 2 else False
            f.color.rgb = r[3] if len(r) > 3 else INK
            _set_ea(run)
    return sp


def line(slide, x1, y1, x2, y2, color=GRAY, width=1.0):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.shadow.inherit = False
    return conn


def arrow(slide, x1, y1, x2, y2, color=GRAY, width=2.0, dash=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.shadow.inherit = False
    ln = conn.line._get_or_add_ln()
    if dash:
        d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.append(d)
    tail = ln.makeelement(qn('a:tailEnd'),
                          {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return conn


def chip(slide, x, y, w, h, txt, fill, color=WHITE, size=14, bold=True):
    sp = box(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.5)
    text_in(sp, [{'runs': [(txt, size, bold, color)]}])
    return sp


def header(slide, title, kicker=None):
    box(slide, 0, 0, 13.333, 1.0, fill=NAVY)
    box(slide, 0.62, 0.33, 0.11, 0.36, fill=GOLD)
    text(slide, 0.9, 0.16, 10.5, 0.68,
         [{'runs': [(title, 28, True, WHITE)]}], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        text(slide, 8.7, 0.3, 4.0, 0.4,
             [{'runs': [(kicker, 12, False, RGBColor(0xB9, 0xC4, 0xDD))],
               'align': PP_ALIGN.RIGHT}], anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, idx):
    text(slide, 0.62, 7.12, 6, 0.3,
         [{'runs': [('Gomoku Formalization · 项目阶段汇报（穿插版）', 9, False, GRAY)]}])
    text(slide, 11.9, 7.12, 0.85, 0.3,
         [{'runs': [(f'{idx} / {TOTAL}', 9, False, GRAY)],
           'align': PP_ALIGN.RIGHT}])


def note_band(slide, y, txt, fill=GOLD_L, color=GOLD_D, size=15, h=0.75,
              bold=True):
    sp = box(slide, 0.62, y, 12.09, h, fill=fill,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
    text_in(sp, [{'runs': [(txt, size, bold, color)]}])
    return sp


def card(slide, x, y, w, h, title, bullets, tcolor=BLUE, title_size=19,
         body_size=15.5, body_y=2.6, step=0.78, fill=LIGHT):
    box(slide, x, y, w, h, fill=fill,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    text(slide, x + 0.33, y + 0.25, w - 0.6, 0.5,
         [{'runs': [(title, title_size, True, tcolor)]}])
    for i, b in enumerate(bullets):
        box(slide, x + 0.35, body_y + i * step + 0.06, 0.15, 0.15, fill=GOLD)
        text(slide, x + 0.66, body_y + i * step, w - 0.95, 0.72,
             [{'runs': [(b, body_size, False, INK)], 'line_spacing': 1.15}])


def issue_card(slide, x, y, w, h, title, desc):
    box(slide, x, y, w, h, fill=RED_L,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    box(slide, x, y, 0.14, h, fill=RED)
    text(slide, x + 0.45, y + 0.28, w - 0.8, 0.55,
         [{'runs': [(title, 19, True, RED)]}])
    text(slide, x + 0.45, y + 1.0, w - 0.8, 1.0,
         [{'runs': [(desc, 15, False, INK)], 'line_spacing': 1.25}])


def three_positions_slide(page, title, kicker, triple, intro):
    """One slide showing three draw positions with board images."""
    s = new_slide()
    header(s, title, kicker)
    text(s, 0.62, 1.18, 12.0, 0.4,
         [{'runs': [(intro, 14.5, False, GRAY)]}])
    xs = [0.62, 4.85, 9.08]
    for i, (lean, seed, png) in enumerate(triple):
        x = xs[i]
        img = os.path.join(BOARD_DIR, png)
        if os.path.exists(img):
            s.shapes.add_picture(img, Inches(x), Inches(1.7), height=Inches(4.3))
        box(s, x, 6.15, 3.63, 0.62, fill=GREEN_L,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
        text(s, x + 0.15, 6.22, 3.35, 0.5,
             [{'runs': [(f'{lean} · {seed}', 12.5, True, GREEN)]}],
             anchor=MSO_ANCHOR.MIDDLE)
    footer(s, page)
    return s


# ================================================================
# Slide 1 · Cover
# ================================================================
s = new_slide(NAVY)
box(s, 0, 7.28, 13.333, 0.22, fill=GOLD)
text(s, 0.9, 1.1, 11.5, 0.45,
     [{'runs': [('GOMOKU FORMALIZATION · PROJECT REPORT', 14, True, GOLD)]}])
text(s, 0.9, 1.8, 11.5, 1.3,
     [{'runs': [('五子棋形式化项目 · 阶段汇报', 46, True, WHITE)]}])
text(s, 0.9, 3.25, 11.0, 0.6,
     [{'runs': [('双模式搜索器 · 两种证书 · 9 个中盘和棋局面', 20, False,
                 RGBColor(0xCF, 0xD8, 0xE8))]}])
box(s, 0.95, 4.15, 2.6, 0.05, fill=GOLD)
text(s, 0.9, 4.4, 11.5, 0.55,
     [{'runs': [('不可信搜索器找候选  ·  Lean 检查器出定理', 18, False,
                 RGBColor(0xF2, 0xCE, 0x94))]}])
text(s, 0.9, 6.3, 11.5, 0.4,
     [{'runs': [('Lean 4 · C++17 DFPN/VCF/DefenseSearcher · CompactCertificate / '
                 'DefenseCertificate', 12.5, False, RGBColor(0x8E, 0x9B, 0xB8))]}])

# ================================================================
# Slide 2 · TOC (threaded)
# ================================================================
s = new_slide()
header(s, '目录', 'Threaded Contents')
toc = [('01', '项目与搜索语义', '两种目标：强制胜 与 防守证明'),
       ('02', '双模式搜索器', 'DFPN/VCF 与 DefenseSearcher'),
       ('03', '证书与可信边界', '两种证书协议 · Lean 检查器'),
       ('04', '实战：9 个中盘和棋局面', '搜索器产出 → 实例 → 验证链'),
       ('05', '不足与总结', '空棋盘目标与工程缺口')]
for i, (num, t1, t2) in enumerate(toc):
    y = 1.32 + i * 1.18
    box(s, 1.15, y, 11.0, 1.0, fill=LIGHT,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    box(s, 1.15, y, 1.18, 1.0, fill=NAVY)
    text(s, 1.15, y, 1.18, 1.0,
         [{'runs': [(num, 28, True, GOLD)], 'align': PP_ALIGN.CENTER}],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, 2.7, y + 0.1, 9.2, 0.5,
         [{'runs': [(t1, 22, True, NAVY)]}])
    text(s, 2.7, y + 0.58, 9.2, 0.4,
         [{'runs': [(t2, 14, False, GRAY)]}])
footer(s, 2)

# ================================================================
# Slide 3 · Project overview
# ================================================================
s = new_slide()
header(s, '项目是什么', 'Overview')
rules = ['棋盘 7 × 7', '五连即胜', '黑棋先手', '无禁手 · 长连也算胜']
for i, r in enumerate(rules):
    chip(s, 0.62 + i * 3.13, 1.45, 2.93, 0.72, r, NAVY, size=16)
text(s, 0.62, 2.55, 12.0, 0.55,
     [{'runs': [('项目做的三件事', 20, True, NAVY)]}])
cards3 = [('形式化规则', '用 Lean 4 定义棋盘、走法与胜负'),
          ('双模式搜索', '强制胜搜索 + 防守证明搜索'),
          ('验证结论', '一切结论都要经过 Lean 检查器')]
for i, (t, d) in enumerate(cards3):
    x = 0.62 + i * 4.13
    box(s, x, 3.25, 3.93, 2.6, fill=LIGHT,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    box(s, x, 3.25, 3.93, 0.72, fill=BLUE)
    text(s, x, 3.25, 3.93, 0.72,
         [{'runs': [(t, 20, True, WHITE)], 'align': PP_ALIGN.CENTER}],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.4, 4.3, 3.2, 1.3,
         [{'runs': [(d, 16.5, False, INK)], 'align': PP_ALIGN.CENTER,
           'line_spacing': 1.3}])
footer(s, 3)

# ================================================================
# Slide 4 · Why a searcher
# ================================================================
s = new_slide()
header(s, '为什么需要搜索器', 'Motivation')
cards = [
    ('① 博弈树太大', '49 个点，手工推演不可能', NAVY),
    ('② 程序会犯错', '搜索结论不能直接当定理', ORANGE := RGBColor(0xD9, 0x7B, 0x29)),
    ('③ 于是分工', '搜索器“找候选”，Lean“出定理”', GREEN),
]
for i, (t1, t2, c) in enumerate(cards):
    x = 0.62 + i * 4.13
    box(s, x, 1.7, 3.93, 2.85, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        radius=0.08)
    text(s, x + 0.4, 2.15, 3.2, 0.6, [{'runs': [(t1, 22, True, WHITE)]}])
    text(s, x + 0.4, 2.95, 3.2, 1.2,
         [{'runs': [(t2, 17, False, WHITE)], 'line_spacing': 1.3}])
label = [('搜索器', '找候选', GREEN), ('证书', '候选数据', ORANGE),
         ('Lean 检查器', '出定理', BLUE)]
x = 2.0
for i, (t1, t2, c) in enumerate(label):
    sp = box(s, x, 5.25, 2.9, 1.1, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.15)
    text_in(sp, [{'runs': [(t1, 18, True, WHITE)]},
                 {'runs': [(t2, 13, False, WHITE)], 'space_after': 2}])
    if i < 2:
        arrow(s, x + 2.9, 5.8, x + 3.75, 5.8, color=GRAY, width=2.4)
    x += 3.75
footer(s, 4)

# ================================================================
# Slide 5 · Two search semantics (force-win vs defense)  [threaded]
# ================================================================
s = new_slide()
header(s, '搜索语义：两种目标', 'Force-Win vs Defense')
# Left: force-win
sp = box(s, 0.62, 1.5, 5.95, 4.6, fill=LIGHT,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
box(s, 0.62, 1.5, 5.95, 0.72, fill=BLUE)
text(s, 0.62, 1.5, 5.95, 0.72,
     [{'runs': [('目标 A · 强制胜 CanForceWin', 19, True, WHITE)],
       'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
bulletsA = [
    ('我方回合 = OR：找到一个能赢的着法即可', GREEN),
    ('对手回合 = AND：每个合法应手都要有胜子树', BLUE),
    ('叶子 = 真实胜局，由检查器重算', GOLD_D),
    ('和棋 / 输棋 / 搜索耗尽 ≠ 赢', RED),
    ('产出：CompactCertificate', ORANGE := RGBColor(0xD9, 0x7B, 0x29)),
]
for i, (t, c) in enumerate(bulletsA):
    y = 2.5 + i * 0.72
    box(s, 0.98, y + 0.05, 0.2, 0.2, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        radius=0.5)
    text(s, 1.4, y, 5.0, 0.6, [{'runs': [(t, 15, True, INK)]}],
         anchor=MSO_ANCHOR.MIDDLE)
# Right: defense
sp = box(s, 6.77, 1.5, 5.95, 4.6, fill=BLUE_L,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
box(s, 6.77, 1.5, 5.95, 0.72, fill=GREEN)
text(s, 6.77, 1.5, 5.95, 0.72,
     [{'runs': [('目标 B · 防守证明 CanPreventWin', 19, True, WHITE)],
       'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
bulletsB = [
    ('防守方回合：选一个保持防守的着法即可', GREEN),
    ('攻击方回合：覆盖全部合法应手，漏一个即拒', BLUE),
    ('叶子 = 防守方胜 或 和棋', GOLD_D),
    ('攻击方胜的终局没有任何闭合路径', RED),
    ('产出：DefenseCertificate', ORANGE),
]
for i, (t, c) in enumerate(bulletsB):
    y = 2.5 + i * 0.72
    box(s, 7.13, y + 0.05, 0.2, 0.2, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        radius=0.5)
    text(s, 7.55, y, 5.0, 0.6, [{'runs': [(t, 15, True, INK)]}],
         anchor=MSO_ANCHOR.MIDDLE)
note_band(s, 6.35, '两种语义都只生成候选数据：搜索器可以犯错，检查器不放错。')
footer(s, 5)

# ================================================================
# Slide 6 · Two-mode searcher architecture
# ================================================================
s = new_slide()
header(s, '搜索器整体架构：双模式', 'Architecture · C++17')
sp = box(s, 0.62, 1.3, 12.09, 0.55, fill=NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
text_in(sp, [{'runs': [('模式 A · 强制胜搜索', 15, True, WHITE)]}])
stagesA = [('输入局面', '7×7 + 轮到谁', NAVY2),
           ('DFPN 主搜索', '证明数优先', BLUE),
           ('VCF 快攻', '连续冲四检测', GOLD_D),
           ('候选证书', 'CompactCertificate', ORANGE),
           ('导出 Lean', 'def certificate', GREEN)]
x = 0.62
for i, (t1, t2, c) in enumerate(stagesA):
    sp = box(s, x, 1.95, 2.16, 1.1, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.12)
    text_in(sp, [{'runs': [(t1, 15, True, WHITE)]},
                 {'runs': [(t2, 11.5, False, WHITE)]}])
    if i < 4:
        arrow(s, x + 2.16, 2.5, x + 2.58, 2.5, color=GRAY, width=2.0)
    x += 2.58
sp = box(s, 0.62, 3.3, 12.09, 0.55, fill=NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
text_in(sp, [{'runs': [('模式 B · 防守证明  （--prove prevent-black-win | '
                        'prevent-white-win）', 15, True, WHITE)]}])
stagesB = [('输入局面', '--input / --root empty', NAVY2),
           ('DefenseSearcher', '完整 AND/OR 搜索', BLUE),
           ('严格传播', 'found / refuted / unknown', GOLD_D),
           ('候选证书', 'DefenseCertificate', ORANGE),
           ('导出 Lean', 'def certificate', GREEN)]
x = 0.62
for i, (t1, t2, c) in enumerate(stagesB):
    sp = box(s, x, 3.95, 2.16, 1.1, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.12)
    text_in(sp, [{'runs': [(t1, 15, True, WHITE)]},
                 {'runs': [(t2, 11.5, False, WHITE)]}])
    if i < 4:
        arrow(s, x + 2.16, 4.5, x + 2.58, 4.5, color=GRAY, width=2.0)
    x += 2.58
chips = ['置换表 / Zobrist 局面键', '迭代加深', '资源上限', '证书预检（仅诊断）']
for i, c in enumerate(chips):
    chip(s, 0.62 + i * 3.13, 5.4, 2.93, 0.6, c, LIGHT, color=INK, size=13)
note_band(s, 6.25, '两条流水线都只产生“候选数据”：最终由 Lean 检查器验证。',
          h=0.55, size=14)
footer(s, 6)

# ================================================================
# Slide 7 · DFPN
# ================================================================
s = new_slide()
header(s, '核心算法 ①：DFPN 证明数搜索', 'DFPN')
cards = [
    ('证明数 pn', BLUE,
     ['证明必胜还需展开的节点数', 'pn 越小 → 越接近证明']),
    ('反证数 dn', RED,
     ['证明必败还需展开的节点数', 'dn 越小 → 分支越该放弃']),
    ('选择策略', GREEN,
     ['优先展开 pn / dn 最小的节点', '深度优先，内存友好',
      '输出证明路径 → 便于生成证书']),
]
for i, (t, c, lines) in enumerate(cards):
    x = 0.62 + i * 4.13
    box(s, x, 1.7, 3.93, 0.75, fill=c)
    text(s, x, 1.7, 3.93, 0.75,
         [{'runs': [(t, 20, True, WHITE)], 'align': PP_ALIGN.CENTER}],
         anchor=MSO_ANCHOR.MIDDLE)
    box(s, x, 2.45, 3.93, 2.9, fill=LIGHT,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    for j, l in enumerate(lines):
        box(s, x + 0.35, 2.85 + j * 0.8, 0.15, 0.15, fill=GOLD)
        text(s, x + 0.66, 2.78 + j * 0.8, 3.05, 0.75,
             [{'runs': [(l, 15.5, False, INK)], 'line_spacing': 1.15}])
note_band(s, 5.7, '把预算花在“最可能证明成功”的分支上；阈值可调：'
                  'maxDepth / maxNodes / maxVcfDepth')
footer(s, 7)

# ================================================================
# Slide 8 · VCF
# ================================================================
s = new_slide()
header(s, '核心算法 ②：VCF 连续冲四', 'VCF')
steps = ['冲四', '对手被迫防守', '再冲四', '对手被迫防守', '成五 ✓']
x = 0.62
for i, t in enumerate(steps):
    c = GREEN if t == '成五 ✓' else NAVY2
    sp = box(s, x, 1.7, 1.75, 0.8, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.2)
    text_in(sp, [{'runs': [(t, 14.5, True, WHITE)]}])
    if i < 4:
        arrow(s, x + 1.75, 2.1, x + 2.31, 2.1, color=GRAY, width=2.2)
    x += 2.31
text(s, 0.62, 2.85, 12.1, 0.5,
     [{'runs': [('冲四：再下一子即成五连，对手必须立刻防守 —— '
                 '连续冲四直到成五，就是一条杀棋路线。', 15.5, False, GRAY)]}])
card(s, 0.62, 3.6, 5.9, 2.9, '为什么有用',
     ['攻击端：快速找到杀棋', '防守端：应手骤减，搜索变小',
      '结果仍要交给检查器复检'],
     tcolor=BLUE, body_size=15.5, body_y=4.4, step=0.72, fill=LIGHT)
card(s, 6.82, 3.6, 5.9, 2.9, '工程要点',
     ['独立预算：maxVcfDepth / maxVcfNodes', '与 DFPN 配合：先试快攻',
      '攻不下来再回到全局面搜索'],
     tcolor=BLUE, body_size=15.5, body_y=4.4, step=0.72, fill=BLUE_L)
footer(s, 8)

# ================================================================
# Slide 9 · Optimizations
# ================================================================
s = new_slide()
header(s, '加速技巧：让搜索更快', 'Optimizations')
card(s, 0.62, 1.7, 5.9, 3.75, '威胁排序 Threat Ordering',
     ['按威胁从高到低先搜', '成五 > 活四 > 冲四 > 活三 …',
      '更快找到杀棋', 'Lean 已证明留下的着法合法'],
     tcolor=BLUE, body_size=16, body_y=2.75, step=0.72, fill=LIGHT)
card(s, 6.82, 1.7, 5.9, 3.75, '强制防守剪枝 Forced-Move Pruning',
     ['对手被逼时合法应手很少', '只展开必要应手',
      '大幅缩小 AND 分支', '配置：useForcedMovePruning'],
     tcolor=BLUE, body_size=16, body_y=2.75, step=0.72, fill=BLUE_L)
note_band(s, 5.75, '两者只影响“找候选”的速度，不改变可信边界 —— '
                   '结果依旧要过 Lean 检查器。')
footer(s, 9)

# ================================================================
# Slide 10 · Resource control & status
# ================================================================
s = new_slide()
header(s, '资源控制：搜不到 ≠ 必败', 'Honest Status')
chips = ['深度上限 maxDepth', '节点预算 maxNodes', '置换表上限', '证书节点上限']
for i, c in enumerate(chips):
    chip(s, 0.62 + i * 3.13, 1.4, 2.93, 0.68, c, NAVY, size=14)
rows = [
    ('found', '找到候选', '交给 Lean 检查器', GREEN),
    ('refuted', '防守被攻破', '无证书（证明不存在防守）', RED),
    ('unknown / 资源耗尽', '搜索被截断', '无证书 · 不做结论', RED),
]
y = 2.5
for name, desc, act, c in rows:
    box(s, 0.62, y, 3.4, 0.78, fill=(GREEN_L if c == GREEN else RED_L),
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
    text(s, 0.62, y, 3.4, 0.78,
         [{'runs': [(name, 15, True, (GREEN if c == GREEN else RED))],
           'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 4.35, y, 3.9, 0.78, [{'runs': [(desc, 16, False, INK)]}],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, 8.5, y, 4.2, 0.78, [{'runs': [(act, 16, False, GRAY)]}],
         anchor=MSO_ANCHOR.MIDDLE)
    y += 0.92
note_band(s, 5.6, '有限搜索失败只表示“本次没找到证据”，绝不是和棋或必败证明；'
                  '提交时须如实报告状态。')
footer(s, 10)

# ================================================================
# Slide 11 · Two certificate protocols
# ================================================================
s = new_slide()
header(s, '输出协议：两种证书', 'Output · v1')
sp = box(s, 0.62, 1.3, 12.09, 0.5, fill=NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
text_in(sp, [{'runs': [('模式 A · CompactCertificate（强制胜）', 14.5, True, WHITE)]}])
cardsA = [
    ('terminal', '胜利叶子', GREEN,
     ['局面真实获胜', '检查器重算胜负', '不信任任何标签']),
    ('proverMove', '目标方回合', BLUE,
     ['选一个合法着法', '指向一个子节点', '子局面 = play s m']),
    ('opponentMoves', '对手回合', ORANGE,
     ['列出全部合法应手', '每个应手一个子节点', '漏一个 → 整张被拒']),
]
for i, (name, sub, c, lines) in enumerate(cardsA):
    x = 0.62 + i * 4.13
    box(s, x, 1.9, 3.93, 0.58, fill=c)
    text(s, x + 0.3, 1.9, 2.3, 0.58, [{'runs': [(name, 17, True, WHITE)]}],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 2.0, 1.9, 1.9, 0.58,
         [{'runs': [(sub, 12.5, False, WHITE)], 'align': PP_ALIGN.RIGHT}],
         anchor=MSO_ANCHOR.MIDDLE)
    box(s, x, 2.48, 3.93, 1.45, fill=LIGHT,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    for j, l in enumerate(lines):
        box(s, x + 0.35, 2.72 + j * 0.4, 0.15, 0.15, fill=GOLD)
        text(s, x + 0.66, 2.65 + j * 0.4, 3.05, 0.4,
             [{'runs': [(l, 13, False, INK)]}])
sp = box(s, 0.62, 4.15, 12.09, 0.5, fill=NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
text_in(sp, [{'runs': [('模式 B · DefenseCertificate（防守证明）', 14.5, True, WHITE)]}])
cardsB = [
    ('terminal', '终局叶子', GREEN,
     ['结果为防守方胜或和棋', '攻击方胜的叶子被拒绝']),
    ('defenderMove', '防守方回合', BLUE,
     ['选一个保持防守的着法', '指向一个子节点']),
    ('attackerMoves', '攻击方回合', ORANGE,
     ['覆盖全部合法应手', '漏一个 → 整张被拒']),
]
for i, (name, sub, c, lines) in enumerate(cardsB):
    x = 0.62 + i * 4.13
    box(s, x, 4.75, 3.93, 0.58, fill=c)
    text(s, x + 0.3, 4.75, 2.3, 0.58, [{'runs': [(name, 17, True, WHITE)]}],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 2.0, 4.75, 1.9, 0.58,
         [{'runs': [(sub, 12.5, False, WHITE)], 'align': PP_ALIGN.RIGHT}],
         anchor=MSO_ANCHOR.MIDDLE)
    box(s, x, 5.33, 3.93, 1.15, fill=LIGHT,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    for j, l in enumerate(lines):
        box(s, x + 0.35, 5.52 + j * 0.4, 0.15, 0.15, fill=GOLD)
        text(s, x + 0.66, 5.45 + j * 0.4, 3.05, 0.4,
             [{'runs': [(l, 13, False, INK)]}])
note_band(s, 6.6, '编号规则：数组下标即编号，parent < child；两种证书都推荐先序布局。',
          h=0.4, size=12.5)
footer(s, 11)

# ================================================================
# Slide 12 · Trust boundary
# ================================================================
s = new_slide()
header(s, '设计核心：可信边界', 'Trust Boundary')
box(s, 0.62, 1.75, 5.15, 3.9, fill=RED_L,
    shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
text(s, 0.95, 2.05, 4.5, 0.55,
     [{'runs': [('不可信：搜索器输出', 19, True, RED)]}])
for i, t in enumerate(['C++ 搜索器（DFPN / VCF / DefenseSearcher）', '启发式与剪枝',
                       '置换表与各种缓存', 'C++ 预检（仅诊断）']):
    box(s, 0.98, 2.8 + i * 0.68, 0.16, 0.16, fill=RED)
    text(s, 1.3, 2.73 + i * 0.68, 4.3, 0.62,
         [{'runs': [(t, 15.5, False, INK)]}])
box(s, 7.56, 1.75, 5.15, 3.9, fill=GREEN_L,
    shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
text(s, 7.89, 2.05, 4.5, 0.55,
     [{'runs': [('可信：Lean 检查器', 19, True, GREEN)]}])
for i, t in enumerate(['checkCertificate / checkDefenseCertificateAt', '重算终局胜负',
                       '验证合法性与子局面', '攻击方全应手覆盖']):
    box(s, 7.92, 2.8 + i * 0.68, 0.16, 0.16, fill=GREEN)
    text(s, 8.24, 2.73 + i * 0.68, 4.3, 0.62,
         [{'runs': [(t, 15.5, False, INK)]}])
arrow(s, 5.77, 3.7, 7.56, 3.7, color=GRAY, width=2.6)
text(s, 5.68, 2.9, 2.0, 0.6,
     [{'runs': [('重新检查', 14, True, GRAY)], 'align': PP_ALIGN.CENTER}])
text(s, 0.95, 5.35, 4.5, 0.4,
     [{'runs': [('→ 只生成候选，可以出错', 14, True, RED)]}])
text(s, 7.89, 5.35, 4.6, 0.4,
     [{'runs': [('→ 通过 → CanForceWin / CanPreventWin 定理', 14, True, GREEN)]}])
note_band(s, 6.05, '搜索器可以犯错，检查器不放错；通过检查的证书才是定理。')
footer(s, 12)

# ================================================================
# Slides 13-15 · Live examples: 9 draw positions (3 per slide)
# ================================================================
three_positions_slide(
    13, '实战：中盘和棋局面 ①', 'DefenseSearcher Output',
    DRAW_POSITIONS[0:3],
    '9 个机器验证的局面 · 21 黑 + 21 白 · 7 空格 · 轮到黑方。'
    '每个长度 5 窗口都含黑白两子，任何一方都无法成五。')

three_positions_slide(
    14, '实战：中盘和棋局面 ②', 'DefenseSearcher Output',
    DRAW_POSITIONS[3:6],
    '每个局面都由 DefenseSearcher 独立生成白防黑 + 黑防白两张防守证书。')

three_positions_slide(
    15, '实战：中盘和棋局面 ③', 'DefenseSearcher Output',
    DRAW_POSITIONS[6:9],
    '证书只是候选：Lean 检查器逐节点重算通过后，才能组合出 StandardDraw 定理。')

# ================================================================
# Slide 16 · Verification chain (draw example threaded here)
# ================================================================
s = new_slide()
header(s, '验证链：从证书到 StandardDraw', 'Trust Chain')
steps = [
    ('① 构造局面', '每个长度 5 窗口都含黑白两子\n→ 任何一方永远无法成五', NAVY2),
    ('② C++ 找证书', '白防黑 + 黑防白 两张防守证书\n(搜索器不可信, 只出候选)', BLUE),
    ('③ Lean 重算', 'checkDefenseCertificateAt\n逐节点重算局面/轮次/终局/应手', GOLD_D),
    ('④ 组合出定理', '两张证书 + soundness\n⇒ StandardDraw 局面', GREEN),
]
x = 0.62
for i, (t1, t2, c) in enumerate(steps):
    sp = box(s, x, 1.7, 2.85, 1.55, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.1)
    text_in(sp, [{'runs': [(t1, 17, True, WHITE)]},
                 {'runs': [(t2, 11.5, False, WHITE)]}])
    if i < 3:
        conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x + 2.85), Inches(2.47),
                                      Inches(x + 3.28), Inches(2.47))
        conn.line.color.rgb = GRAY
        conn.line.width = Pt(2.0)
    x += 3.28
text(s, 0.62, 3.75, 12.1, 0.55,
     [{'runs': [('每个局面独立验证 4 项断言 (全部 native_decide, 不依赖 C++ 结果):',
                 16, True, NAVY)]}])
checks = [
    ('白防黑证书 checkDefenseCertificateAt = true', '9 / 9 通过'),
    ('黑防白证书 checkDefenseCertificateAt = true', '9 / 9 通过'),
    ('局面形状: 21黑 + 21白 + 7空 + 轮到黑方', '9 / 9 通过'),
    ('60 个长度 5 窗口全部双色', '9 / 9 通过'),
    ('StandardDraw <position> 由两张证书组合', '9 / 9 通过'),
]
y = 4.4
for i, (t, r) in enumerate(checks):
    box(s, 0.62, y, 8.6, 0.42, fill=LIGHT,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    text(s, 0.95, y, 8.2, 0.42, [{'runs': [(t, 14, False, INK)]}],
         anchor=MSO_ANCHOR.MIDDLE)
    chip = box(s, 9.5, y, 3.0, 0.42, fill=GREEN_L,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    text_in(chip, [{'runs': [(r, 13.5, True, GREEN)]}])
    y += 0.5
note = box(s, 0.62, 6.5, 12.09, 0.5, fill=GOLD_L,
           shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.2)
text_in(note, [{'runs': [('搜索器可以犯错, 检查器不放错; 通过检查的证书才是定理。',
                          14, True, GOLD_D)]}])
footer(s, 16)

# ================================================================
# Slide 17 · Summary table
# ================================================================
s = new_slide()
header(s, '全部 9 个局面: 汇总', 'Summary')
col_x = [0.62, 3.1, 5.2, 7.0, 8.9, 10.9]
col_w = [2.48, 2.1, 1.8, 1.9, 2.0, 1.6]
heads = ['Lean 定理文件', '种子', '白防证书', '黑防证书', 'StandardDraw', '窗口双色']
for i, h in enumerate(heads):
    box(s, col_x[i], 1.5, col_w[i], 0.55, fill=NAVY)
    text(s, col_x[i], 1.5, col_w[i], 0.55,
         [{'runs': [(h, 14, True, WHITE)], 'align': PP_ALIGN.CENTER}],
         anchor=MSO_ANCHOR.MIDDLE)
y = 2.15
for idx, (lean, seed, png) in enumerate(DRAW_POSITIONS):
    fill = LIGHT if idx % 2 == 0 else WHITE
    for i, val in enumerate([lean, seed, '✓', '✓', '✓ 定理', '✓']):
        box(s, col_x[i], y, col_w[i], 0.52, fill=fill)
        color = GREEN if val == '✓' or val == '✓ 定理' else INK
        text(s, col_x[i], y, col_w[i], 0.52,
             [{'runs': [(val, 13, (val != lean), color)],
               'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.52
band = box(s, 0.62, 6.35, 12.09, 0.6, fill=GREEN_L,
           shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
text_in(band, [{'runs': [('9 / 9 全部通过: 每张证书经 Lean 检查器重算, 每个局面都证明了 '
                          'StandardDraw。', 15, True, GREEN)]}])
footer(s, 17)

# ================================================================
# Slide 18 · Gaps 1
# ================================================================
s = new_slide()
header(s, '不足 ①：功能缺口', 'Not Done Well · Function')
issues = [
    ('空棋盘开局求解 · 未完成',
     '没有从 7×7 空棋盘开始的必胜证书，不能宣称 initial_black_wins'),
    ('空棋盘和棋 · 未完成',
     '中盘 9 局面的和棋已验证；空棋盘 StandardDraw 仍是开放目标'),
    ('缓存 · 未接入',
     'SearchMemo / EngineMemo 已定义，但还没有接入引擎默认搜索流程'),
    ('性能 · 未验证',
     '目前只有一步胜、双应手、开四 VCF 与中盘和棋等示例，规模没有跑起来'),
]
for i, (t, d) in enumerate(issues):
    x = 0.62 + (i % 2) * 6.2
    y = 1.7 + (i // 2) * 2.35
    issue_card(s, x, y, 5.9, 2.1, t, d)
note_band(s, 6.35, '这些缺口意味着：项目目前给不出“空棋盘黑棋先手必胜”或“空棋盘和棋”的结论。',
          fill=RED_L, color=RED)
footer(s, 18)

# ================================================================
# Slide 19 · Gaps 2
# ================================================================
s = new_slide()
header(s, '不足 ②：工程与验证', 'Not Done Well · Engineering')
issues = [
    ('参数没有系统调优',
     'maxDepth / maxVcfDepth / maxNodes 用的是保守默认值，没做过参数扫描'),
    ('缺少基准数据',
     '证书大小、Lean 检查耗时、缓存命中率都没有系统记录和对比'),
    ('导出→复检链路粗糙',
     'C++ 导出 Lean 源码后要手工跑检查，自动化程度不高'),
    ('投入偏重正确性',
     '时间主要花在规则与检查器上，搜索规模与性能投入不足'),
]
for i, (t, d) in enumerate(issues):
    x = 0.62 + (i % 2) * 6.2
    y = 1.7 + (i // 2) * 2.35
    issue_card(s, x, y, 5.9, 2.1, t, d)
note_band(s, 6.35, '一句话：框架完整、正确性优先，但离完整 7×7 求解还有明显距离。',
          fill=RED_L, color=RED)
footer(s, 19)

# ================================================================
# Slide 20 · Closing
# ================================================================
s = new_slide(NAVY)
box(s, 0, 7.28, 13.333, 0.22, fill=GOLD)
text(s, 0.9, 1.8, 11.5, 1.3,
     [{'runs': [('谢谢', 60, True, WHITE)], 'align': PP_ALIGN.CENTER}])
box(s, 6.17, 3.35, 1.0, 0.05, fill=GOLD)
text(s, 0.9, 3.7, 11.5, 0.7,
     [{'runs': [('搜索器负责“找”，Lean 负责“证”', 24, False,
                 RGBColor(0xF2, 0xCE, 0x94))], 'align': PP_ALIGN.CENTER}])
text(s, 0.9, 4.7, 11.5, 0.55,
     [{'runs': [('已做：规则 · 双模式搜索器 · 两种证书 · 9 个中盘和棋局面', 16, False,
                 RGBColor(0xCF, 0xD8, 0xE8))], 'align': PP_ALIGN.CENTER}])
text(s, 0.9, 5.3, 11.5, 0.55,
     [{'runs': [('未做：空棋盘开局求解 · 空棋盘和棋 · 大规模验证', 16, False,
                 RGBColor(0xCF, 0xD8, 0xE8))], 'align': PP_ALIGN.CENTER}])

# ---------------- speaker notes ----------------
NOTES = [
    '这一版把 9 个和棋局面穿插进技术主线，而不是单独一章：先讲两种搜索语义，'
    '再讲双模式搜索器、两种证书协议与可信边界，然后用三页实战局面和验证链收尾。',
    '目录按主题流动：语义、双模式、证书与边界、实战局面、不足与总结。',
    '项目用 Lean 4 形式化固定规则的五子棋：7×7、五连即胜、黑先、无禁手。',
    '为什么需要搜索器：博弈树太大，程序会犯错，所以分工：搜索器找候选，Lean 出定理。',
    '这一页是穿插的关键：两种目标并排。目标 A 强制胜：我方回合 OR、对手回合 AND；'
    '目标 B 防守证明：防守方选一步、攻击方全应手覆盖，叶子只能是防守方胜或和棋。'
    '和棋局面就是目标 B 的产物。',
    '双模式架构：模式 A 是 DFPN + VCF 强制胜，模式 B 是 DefenseSearcher 完整 AND/OR '
    '防守证明，状态严格区分 found、refuted、unknown。两条流水线都只出候选。',
    'DFPN 用证明数 pn 与反证数 dn 优先展开最有希望的分支。',
    'VCF 连续冲四快攻，攻不下来退回 DFPN。',
    '威胁排序与强制防守剪枝只影响速度，不影响正确性。',
    '资源控制：found 才出证书；refuted 是防守被攻破；unknown 是搜索被截断，'
    '绝不能写成和棋或必败。',
    '两种证书协议：CompactCertificate 服务强制胜，DefenseCertificate 服务防守证明；'
    'DefenseCertificate 的终局叶子只接受防守方胜或和棋。',
    '可信边界：左边全部不可信只出候选，右边 Lean 检查器逐节点重算，'
    '通过才得到 CanForceWin 或 CanPreventWin 定理。',
    '实战第一页：前三个和棋局面。每个局面每个长度 5 窗口都含黑白两子。',
    '实战第二页：中间三个局面，都由 DefenseSearcher 生成两张防守证书。',
    '实战第三页：最后三个局面。证书只是候选，Lean 检查通过才是定理。',
    '验证链：构造局面、C++ 找证书、Lean 重算、组合出 StandardDraw；五项断言 9/9 通过。',
    '汇总表：九个局面的两张证书全部通过检查，每个都证明了 StandardDraw。',
    '不足一：空棋盘开局求解和空棋盘和棋都未完成，目前只有中盘局面。',
    '不足二：参数没调优、缺基准数据、导出复检链路粗糙。',
    '总结：已做规则、双模式搜索器、两种证书与 9 个中盘和棋；未做空棋盘目标。谢谢。',
]
for slide, note in zip(prs.slides, NOTES):
    slide.notes_slide.notes_text_frame.text = note

prs.core_properties.title = '五子棋形式化项目阶段汇报（穿插版）'
prs.core_properties.author = 'Gomoku Formalization'

OUT = os.path.join(REPO_ROOT, '五子棋项目阶段汇报_穿插版.pptx')
try:
    prs.save(OUT)
    print('saved:', OUT, 'slides =', len(prs.slides._sldIdLst))
except PermissionError:
    OUT = os.path.join(REPO_ROOT, '五子棋项目阶段汇报_穿插版_新版.pptx')
    prs.save(OUT)
    print('original file locked; saved as:', OUT,
          'slides =', len(prs.slides._sldIdLst))

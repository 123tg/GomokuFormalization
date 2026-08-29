# -*- coding: utf-8 -*-
"""Inspect the merged PPT: slide count, pictures per slide, first text, notes."""
from pptx import Presentation

p = Presentation(r'C:\Users\lenovo\Desktop\GomokuFormalization\五子棋项目阶段汇报.pptx')
print('slides:', len(p.slides._sldIdLst))
for i, s in enumerate(p.slides, 1):
    pics = [sh for sh in s.shapes if sh.shape_type == 13]
    texts = []
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            t = sh.text_frame.text.strip().split('\n')[0]
            if t and len(t) > 2:
                texts.append(t[:34])
    note = ''
    try:
        if s.notes_slide.notes_text_frame.text.strip():
            note = ' [notes]'
    except Exception:
        pass
    first = texts[0] if texts else ''
    print('{:2d}: pics={} {!r:38s}{}'.format(i, len(pics), first, note))

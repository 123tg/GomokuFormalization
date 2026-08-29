# -*- coding: utf-8 -*-
"""Generate a report-style PPT for the Gomoku project.
Honest status report: what is done, what is NOT done well.
Big fonts, light content per slide."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------------- palette ----------------
NAVY    = RGBColor(0x1B, 0x2A, 0x4A)
NAVY2   = RGBColor(0x2C, 0x3E, 0x63)
BLUE    = RGBColor(0x3B, 0x59, 0x98)
BLUE_L  = RGBColor(0xE7, 0xEC, 0xF6)
GOLD    = RGBColor(0xE8, 0xA3, 0x3D)
GOLD_D  = RGBColor(0xC9, 0x8A, 0x2D)
GOLD_L  = RGBColor(0xFB, 0xEF, 0xDA)
GREEN   = RGBColor(0x2E, 0x7D, 0x4F)
GREEN_L = RGBColor(0xE4, 0xF1, 0xE9)
RED     = RGBColor(0xB3, 0x40, 0x2F)
RED_L   = RGBColor(0xF8, 0xE8, 0xE4)
ORANGE  = RGBColor(0xD9, 0x7B, 0x29)
INK     = RGBColor(0x2B, 0x2B, 0x2B)
GRAY    = RGBColor(0x6B, 0x74, 0x80)
LIGHT   = RGBColor(0xF1, 0xF3, 0xF7)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
WOOD    = RGBColor(0xF0, 0xD9, 0xB0)
WOOD_L  = RGBColor(0x9A, 0x7B, 0x4F)
FONT    = '微软雅黑'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL = 28

# 9 machine-verified 7×7 draw positions (lean file, seed, board PNG).
BOARD_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                         'LeanProjects-work', 'docs', 'boards')
DRAW_POSITIONS = [
    ('Draw7x7.lean',   'seed 1', 'Draw7x7.png'),
    ('Draw7x7s2.lean', 'seed 2', 'Draw7x7s2.png'),
    ('Draw7x7s3.lean', 'seed 3', 'Draw7x7s3.png'),
    ('Draw7x7s4.lean', 'seed 4', 'Draw7x7s4.png'),
    ('Draw7x7s5.lean', 'seed 5', 'Draw7x7s5.png'),
    ('Draw7x7s6.lean', 'seed 6', 'Draw7x7s6.png'),
    ('Draw7x7s7.lean', 'seed 7', 'Draw7x7s7.png'),
    ('Draw7x7s8.lean', 'seed 8', 'Draw7x7s8.png'),
    ('Draw7x7s9.lean', 'seed 9', 'Draw7x7s9.png'),
]


def new_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def _set_ea(run, name=FONT):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', name)


def text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = para.get('align', align)
        if para.get('space_before') is not None:
            p.space_before = Pt(para['space_before'])
        if para.get('space_after') is not None:
            p.space_after = Pt(para['space_after'])
        if para.get('line_spacing') is not None:
            p.line_spacing = para['line_spacing']
        for r in para['runs']:
            run = p.add_run()
            run.text = r[0]
            f = run.font
            f.name = FONT
            f.size = Pt(r[1])
            f.bold = r[2] if len(r) > 2 else False
            f.color.rgb = r[3] if len(r) > 3 else INK
            _set_ea(run)
    return tb


def box(slide, x, y, w, h, fill=None, line=None, lw=0.75,
        shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    if radius is not None:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def text_in(sp, paras, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = para.get('align', align)
        if para.get('space_after') is not None:
            p.space_after = Pt(para['space_after'])
        for r in para['runs']:
            run = p.add_run()
            run.text = r[0]
            f = run.font
            f.name = FONT
            f.size = Pt(r[1])
            f.bold = r[2] if len(r) > 2 else False
            f.color.rgb = r[3] if len(r) > 3 else INK
            _set_ea(run)
    return sp


def line(slide, x1, y1, x2, y2, color=GRAY, width=1.0):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.shadow.inherit = False
    return conn


def arrow(slide, x1, y1, x2, y2, color=GRAY, width=2.0, dash=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.shadow.inherit = False
    ln = conn.line._get_or_add_ln()
    if dash:
        d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.append(d)
    tail = ln.makeelement(qn('a:tailEnd'),
                          {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return conn


def chip(slide, x, y, w, h, txt, fill, color=WHITE, size=14, bold=True):
    sp = box(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.5)
    text_in(sp, [{'runs': [(txt, size, bold, color)]}])
    return sp


def header(slide, title, kicker=None):
    box(slide, 0, 0, 13.333, 1.0, fill=NAVY)
    box(slide, 0.62, 0.33, 0.11, 0.36, fill=GOLD)
    text(slide, 0.9, 0.16, 10.5, 0.68,
         [{'runs': [(title, 28, True, WHITE)]}], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        text(slide, 8.7, 0.3, 4.0, 0.4,
             [{'runs': [(kicker, 12, False, RGBColor(0xB9, 0xC4, 0xDD))],
               'align': PP_ALIGN.RIGHT}], anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, idx):
    text(slide, 0.62, 7.12, 6, 0.3,
         [{'runs': [('Gomoku Formalization · 项目阶段汇报', 9, False, GRAY)]}])
    text(slide, 11.9, 7.12, 0.85, 0.3,
         [{'runs': [(f'{idx} / {TOTAL}', 9, False, GRAY)],
           'align': PP_ALIGN.RIGHT}])


def note_band(slide, y, txt, fill=GOLD_L, color=GOLD_D, size=15, h=0.75,
              bold=True):
    sp = box(slide, 0.62, y, 12.09, h, fill=fill,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
    text_in(sp, [{'runs': [(txt, size, bold, color)]}])


def draw_board(slide, x, y, size, stones):
    """Draw a small 7x7 gomoku board. stones: {(col,row): 'B'|'W'}"""
    box(slide, x - 0.10, y - 0.10, size + 0.20, size + 0.20,
        fill=WOOD, line=WOOD_L, lw=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        radius=0.08)
    cell = size / 7.0
    for i in range(8):
        off = i * cell
        line(slide, x, y + off, x + size, y + off, color=WOOD_L, width=1.0)
        line(slide, x + off, y, x + off, y + size, color=WOOD_L, width=1.0)
    cx = x + 3 * cell + cell / 2
    cy = y + 3 * cell + cell / 2
    box(slide, cx - 0.035, cy - 0.035, 0.07, 0.07, fill=WOOD_L,
        shape=MSO_SHAPE.OVAL)
    for (col, row), who in stones.items():
        d = cell * 0.80
        sx = x + col * cell + (cell - d) / 2
        sy = y + row * cell + (cell - d) / 2
        if who == 'B':
            box(slide, sx, sy, d, d, fill=RGBColor(0x20, 0x20, 0x20),
                shape=MSO_SHAPE.OVAL)
        else:
            box(slide, sx, sy, d, d, fill=WHITE, line=RGBColor(0xB8, 0xB8, 0xB8),
                lw=1.0, shape=MSO_SHAPE.OVAL)


def card(slide, x, y, w, h, title, bullets, tcolor=BLUE, title_size=19,
         body_size=15.5, body_y=2.6, step=0.78, fill=LIGHT):
    box(slide, x, y, w, h, fill=fill,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    text(slide, x + 0.33, y + 0.25, w - 0.6, 0.5,
         [{'runs': [(title, title_size, True, tcolor)]}])
    for i, b in enumerate(bullets):
        box(slide, x + 0.35, body_y + i * step + 0.06, 0.15, 0.15, fill=GOLD)
        text(slide, x + 0.66, body_y + i * step, w - 0.95, 0.72,
             [{'runs': [(b, body_size, False, INK)], 'line_spacing': 1.15}])


def issue_card(slide, x, y, w, h, title, desc):
    """Red-tinted card for 'not done well' items."""
    box(slide, x, y, w, h, fill=RED_L,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    box(slide, x, y, 0.14, h, fill=RED)
    text(slide, x + 0.45, y + 0.28, w - 0.8, 0.55,
         [{'runs': [(title, 19, True, RED)]}])
    text(slide, x + 0.45, y + 1.0, w - 0.8, 1.0,
         [{'runs': [(desc, 15, False, INK)], 'line_spacing': 1.25}])


# ================================================================
# Slide 1 · Cover
# ================================================================
s = new_slide(NAVY)
box(s, 0, 7.28, 13.333, 0.22, fill=GOLD)
text(s, 0.9, 1.15, 9.5, 0.45,
     [{'runs': [('GOMOKU FORMALIZATION · PROJECT REPORT', 14, True, GOLD)]}])
text(s, 0.9, 1.85, 9.5, 1.3,
     [{'runs': [('五子棋形式化项目 · 阶段汇报', 48, True, WHITE)]}])
text(s, 0.9, 3.3, 9.2, 0.55,
     [{'runs': [('搜索器 · 证书检查 · 现状与不足', 20, False,
                 RGBColor(0xCF, 0xD8, 0xE8))]}])
box(s, 0.95, 4.2, 2.6, 0.05, fill=GOLD)
text(s, 0.9, 4.5, 9.2, 0.55,
     [{'runs': [('不可信搜索器找候选  ·  Lean 检查器出定理', 18, False,
                 RGBColor(0xF2, 0xCE, 0x94))]}])
text(s, 0.9, 6.35, 9.2, 0.4,
     [{'runs': [('Lean 4 · C++17 DFPN/VCF · CompactCertificate 证书检查',
                 13, False, RGBColor(0x8E, 0x9B, 0xB8))]}])
draw_board(s, 9.3, 1.7, 3.4,
           {(1, 3): 'B', (2, 3): 'B', (3, 3): 'B', (4, 3): 'B', (5, 3): 'B',
            (0, 1): 'W', (6, 5): 'W'})

# ================================================================
# Slide 2 · TOC
# ================================================================
s = new_slide()
header(s, '目录', 'Contents')
toc = [('01', '项目概览', '规则 · 架构 · 搜索语义'),
       ('02', '搜索器与验证', '算法 · 证书 · 可信边界'),
       ('03', '已完成成果', '做到什么程度'),
       ('04', '和棋局面验证', '9 个局面 · 机器验证 · 棋盘图'),
       ('05', '不足与反思', '哪些地方还做得不好')]
for i, (num, t1, t2) in enumerate(toc):
    y = 1.32 + i * 1.18
    box(s, 1.15, y, 11.0, 1.18, fill=LIGHT,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    box(s, 1.15, y, 1.18, 1.18, fill=NAVY)
    text(s, 1.15, y, 1.18, 1.18,
         [{'runs': [(num, 30, True, GOLD)], 'align': PP_ALIGN.CENTER}],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, 2.7, y + 0.14, 9.2, 0.55,
         [{'runs': [(t1, 24, True, NAVY)]}])
    text(s, 2.7, y + 0.68, 9.2, 0.45,
         [{'runs': [(t2, 15, False, GRAY)]}])
footer(s, 2)

# ================================================================
# Slide 3 · Project overview
# ================================================================
s = new_slide()
header(s, '项目是什么', 'Overview')
rules = ['棋盘 7 × 7', '五连即胜', '黑棋先手', '无禁手 · 长连也算胜']
for i, r in enumerate(rules):
    chip(s, 0.62 + i * 3.13, 1.45, 2.93, 0.72, r, NAVY, size=16)
text(s, 0.62, 2.55, 12.0, 0.55,
     [{'runs': [('项目做的三件事', 20, True, NAVY)]}])
cards3 = [('形式化规则', '用 Lean 4 定义棋盘、走法与胜负'),
          ('搜索策略', '用搜索器寻找“黑棋必胜策略树”'),
          ('验证结论', '一切结论都要经过 Lean 检查器')]
for i, (t, d) in enumerate(cards3):
    x = 0.62 + i * 4.13
    box(s, x, 3.25, 3.93, 2.6, fill=LIGHT,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    box(s, x, 3.25, 3.93, 0.72, fill=BLUE)
    text(s, x, 3.25, 3.93, 0.72,
         [{'runs': [(t, 20, True, WHITE)], 'align': PP_ALIGN.CENTER}],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.4, 4.3, 3.2, 1.3,
         [{'runs': [(d, 16.5, False, INK)], 'align': PP_ALIGN.CENTER,
           'line_spacing': 1.3}])
footer(s, 3)

# ================================================================
# Slide 4 · Why a searcher
# ================================================================
s = new_slide()
header(s, '为什么需要搜索器', 'Motivation')
cards = [
    ('① 博弈树太大', '49 个点，手工推演不可能', NAVY),
    ('② 程序会犯错', '搜索结论不能直接当定理', ORANGE),
    ('③ 于是分工', '搜索器“找候选”，Lean“出定理”', GREEN),
]
for i, (t1, t2, c) in enumerate(cards):
    x = 0.62 + i * 4.13
    box(s, x, 1.7, 3.93, 2.85, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        radius=0.08)
    text(s, x + 0.4, 2.15, 3.2, 0.6, [{'runs': [(t1, 22, True, WHITE)]}])
    text(s, x + 0.4, 2.95, 3.2, 1.2,
         [{'runs': [(t2, 17, False, WHITE)], 'line_spacing': 1.3}])
label = [('搜索器', '找候选', GREEN), ('证书', '候选数据', ORANGE),
         ('Lean 检查器', '出定理', BLUE)]
x = 2.0
for i, (t1, t2, c) in enumerate(label):
    sp = box(s, x, 5.25, 2.9, 1.1, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.15)
    text_in(sp, [{'runs': [(t1, 18, True, WHITE)]},
                 {'runs': [(t2, 13, False, WHITE)], 'space_after': 2}])
    if i < 2:
        arrow(s, x + 2.9, 5.8, x + 3.75, 5.8, color=GRAY, width=2.4)
    x += 3.75
footer(s, 4)

# ================================================================
# Slide 5 · Force-win tree (AND/OR)
# ================================================================
s = new_slide()
header(s, '搜索语义：必胜策略树', 'AND / OR Tree')
or_sp = box(s, 1.0, 1.95, 3.6, 1.05, fill=GOLD,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
text_in(or_sp, [{'runs': [('我方回合  ·  OR', 18, True, WHITE)]},
                {'runs': [('找一个能赢的着法即可', 13, False, WHITE)]}])
and_sp = box(s, 1.0, 3.95, 3.6, 1.05, fill=BLUE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
text_in(and_sp, [{'runs': [('对手回合  ·  AND', 18, True, WHITE)]},
                 {'runs': [('每个合法应手都要赢', 13, False, WHITE)]}])
arrow(s, 2.8, 3.0, 2.8, 3.95, color=GRAY, width=2.2)
leaves = [('胜局 ✓', GREEN), ('胜局 ✓', GREEN), ('胜局 ✓', GREEN),
          ('漏应手 ✗', RED)]
for i, (t, c) in enumerate(leaves):
    x = 0.78 + i * 1.14
    sp = box(s, x, 5.75, 1.0, 0.72, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.2)
    text_in(sp, [{'runs': [(t, 12.5, True, WHITE)]}])
    arrow(s, 2.8, 5.0, x + 0.5, 5.75,
          color=(GRAY if c != RED else RED), width=1.8, dash=(c == RED))
text(s, 0.78, 6.6, 4.0, 0.45,
     [{'runs': [('对手有几个合法应手，就要展开几个分支', 12.5, False, GRAY)],
       'align': PP_ALIGN.CENTER}])
bx = 5.45
bullets = [
    ('目标方回合：找到一个成功着法即可', GREEN),
    ('对手回合：所有应手都要有胜子树', BLUE),
    ('叶子：真实胜局，由检查器重算', GOLD_D),
    ('和棋 / 输棋 / 搜索耗尽 ≠ 赢', RED),
]
for i, (t, c) in enumerate(bullets):
    y = 2.0 + i * 0.95
    box(s, bx, y + 0.05, 0.2, 0.2, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        radius=0.5)
    text(s, bx + 0.45, y, 7.0, 0.5,
         [{'runs': [(t, 17.5, True, INK)]}], anchor=MSO_ANCHOR.MIDDLE)
text(s, bx, 6.1, 7.3, 0.6,
     [{'runs': [('这就是 ', 15, False, INK),
                ('CanForceWin', 15, True, BLUE),
                (' 的数学含义：搜索器只生成数据，证明由 Lean 决定。',
                 15, False, INK)],
       'line_spacing': 1.2}])
footer(s, 5)

# ================================================================
# Slide 6 · Searcher architecture
# ================================================================
s = new_slide()
header(s, '搜索器整体架构', 'Architecture · C++17')
stages = [('输入局面', '7×7 + 轮到谁', NAVY2),
          ('DFPN 主搜索', '证明数优先', BLUE),
          ('VCF 快攻', '连续冲四检测', GOLD_D),
          ('候选证书', 'CompactCertificate', ORANGE),
          ('导出 Lean 源码', 'def certificate', GREEN)]
x = 0.62
for i, (t1, t2, c) in enumerate(stages):
    sp = box(s, x, 1.85, 2.16, 1.3, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.12)
    text_in(sp, [{'runs': [(t1, 16, True, WHITE)]},
                 {'runs': [(t2, 12, False, WHITE)]}])
    if i < 4:
        arrow(s, x + 2.16, 2.5, x + 2.58, 2.5, color=GRAY, width=2.2)
    x += 2.58
chips = ['置换表 / Zobrist 局面键', '迭代加深', '资源上限', '证书预检（仅诊断）']
for i, c in enumerate(chips):
    chip(s, 0.62 + i * 3.13, 3.85, 2.93, 0.68, c, LIGHT, color=INK, size=14)
note_band(s, 5.05, '整条流水线只产生“候选数据”：任何算法、启发式、缓存'
                   '都不能直接当结论，最终由 Lean 检查器验证。', h=0.8, size=15.5)
footer(s, 6)

# ================================================================
# Slide 7 · DFPN
# ================================================================
s = new_slide()
header(s, '核心算法 ①：DFPN 证明数搜索', 'DFPN')
cards = [
    ('证明数 pn', BLUE,
     ['证明必胜还需展开的节点数', 'pn 越小 → 越接近证明']),
    ('反证数 dn', RED,
     ['证明必败还需展开的节点数', 'dn 越小 → 分支越该放弃']),
    ('选择策略', GREEN,
     ['优先展开 pn / dn 最小的节点', '深度优先，内存友好',
      '输出证明路径 → 便于生成证书']),
]
for i, (t, c, lines) in enumerate(cards):
    x = 0.62 + i * 4.13
    box(s, x, 1.7, 3.93, 0.75, fill=c)
    text(s, x, 1.7, 3.93, 0.75,
         [{'runs': [(t, 20, True, WHITE)], 'align': PP_ALIGN.CENTER}],
         anchor=MSO_ANCHOR.MIDDLE)
    box(s, x, 2.45, 3.93, 2.9, fill=LIGHT,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    for j, l in enumerate(lines):
        box(s, x + 0.35, 2.85 + j * 0.8, 0.15, 0.15, fill=GOLD)
        text(s, x + 0.66, 2.78 + j * 0.8, 3.05, 0.75,
             [{'runs': [(l, 15.5, False, INK)], 'line_spacing': 1.15}])
note_band(s, 5.7, '把预算花在“最可能证明成功”的分支上；阈值可调：'
                  'maxDepth / maxNodes / maxVcfDepth')
footer(s, 7)

# ================================================================
# Slide 8 · VCF
# ================================================================
s = new_slide()
header(s, '核心算法 ②：VCF 连续冲四', 'VCF')
steps = ['冲四', '对手被迫防守', '再冲四', '对手被迫防守', '成五 ✓']
x = 0.62
for i, t in enumerate(steps):
    c = GREEN if t == '成五 ✓' else NAVY2
    sp = box(s, x, 1.7, 1.75, 0.8, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.2)
    text_in(sp, [{'runs': [(t, 14.5, True, WHITE)]}])
    if i < 4:
        arrow(s, x + 1.75, 2.1, x + 2.31, 2.1, color=GRAY, width=2.2)
    x += 2.31
text(s, 0.62, 2.85, 12.1, 0.5,
     [{'runs': [('冲四：再下一子即成五连，对手必须立刻防守 —— '
                 '连续冲四直到成五，就是一条杀棋路线。', 15.5, False, GRAY)]}])
card(s, 0.62, 3.6, 5.9, 2.9, '为什么有用',
     ['攻击端：快速找到杀棋', '防守端：应手骤减，搜索变小',
      '结果仍要交给检查器复检'],
     tcolor=BLUE, body_size=15.5, body_y=4.4, step=0.72, fill=LIGHT)
card(s, 6.82, 3.6, 5.9, 2.9, '工程要点',
     ['独立预算：maxVcfDepth / maxVcfNodes', '与 DFPN 配合：先试快攻',
      '攻不下来再回到全局面搜索'],
     tcolor=BLUE, body_size=15.5, body_y=4.4, step=0.72, fill=BLUE_L)
footer(s, 8)

# ================================================================
# Slide 9 · Threat ordering & forced-move pruning
# ================================================================
s = new_slide()
header(s, '加速技巧：让搜索更快', 'Optimizations')
card(s, 0.62, 1.7, 5.9, 3.75, '威胁排序 Threat Ordering',
     ['按威胁从高到低先搜', '成五 > 活四 > 冲四 > 活三 …',
      '更快找到杀棋', 'Lean 已证明留下的着法合法'],
     tcolor=BLUE, body_size=16, body_y=2.75, step=0.72, fill=LIGHT)
card(s, 6.82, 1.7, 5.9, 3.75, '强制防守剪枝 Forced-Move Pruning',
     ['对手被逼时合法应手很少', '只展开必要应手',
      '大幅缩小 AND 分支', '配置：useForcedMovePruning'],
     tcolor=BLUE, body_size=16, body_y=2.75, step=0.72, fill=BLUE_L)
note_band(s, 5.75, '两者只影响“找候选”的速度，不改变可信边界 —— '
                   '结果依旧要过 Lean 检查器。')
footer(s, 9)

# ================================================================
# Slide 10 · Resource control & status
# ================================================================
s = new_slide()
header(s, '资源控制：搜不到 ≠ 必败', 'Honest Status')
chips = ['深度上限 maxDepth', '节点预算 maxNodes', '置换表上限', '证书节点上限']
for i, c in enumerate(chips):
    chip(s, 0.62 + i * 3.13, 1.4, 2.93, 0.68, c, NAVY, size=14)
rows = [
    ('found', '找到候选', '交给 Lean 检查器', GREEN),
    ('depthLimit', '深度耗尽', '无证书', RED),
    ('nodeLimit', '节点耗尽', '无证书', RED),
    ('tableLimit / certificateLimit', '表上限', '无证书', RED),
]
y = 2.5
for name, desc, act, c in rows:
    box(s, 0.62, y, 3.4, 0.78, fill=(GREEN_L if c == GREEN else RED_L),
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
    text(s, 0.62, y, 3.4, 0.78,
         [{'runs': [(name, 15, True, (GREEN if c == GREEN else RED))],
           'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 4.35, y, 3.9, 0.78, [{'runs': [(desc, 16, False, INK)]}],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, 8.5, y, 4.2, 0.78, [{'runs': [(act, 16, False, GRAY)]}],
         anchor=MSO_ANCHOR.MIDDLE)
    y += 0.92
note_band(s, 6.3, '有限搜索失败只表示“本次没找到证据”，绝不是和棋或必败证明；'
                  '提交时须如实报告状态。')
footer(s, 10)

# ================================================================
# Slide 11 · Certificate format
# ================================================================
s = new_slide()
header(s, '输出协议：CompactCertificate', 'Output')
cards = [
    ('terminal', '胜利叶子', GREEN,
     ['局面真实获胜', '检查器重算胜负', '不信任任何标签']),
    ('proverMove', '目标方回合', BLUE,
     ['选一个合法着法', '指向一个子节点', '子局面 = play s m']),
    ('opponentMoves', '对手回合', ORANGE,
     ['列出全部合法应手', '每个应手一个子节点', '漏一个 → 整张被拒']),
]
for i, (name, sub, c, lines) in enumerate(cards):
    x = 0.62 + i * 4.13
    box(s, x, 1.65, 3.93, 0.78, fill=c)
    text(s, x + 0.3, 1.65, 2.3, 0.78, [{'runs': [(name, 18, True, WHITE)]}],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 2.0, 1.65, 1.9, 0.78,
         [{'runs': [(sub, 13.5, False, WHITE)], 'align': PP_ALIGN.RIGHT}],
         anchor=MSO_ANCHOR.MIDDLE)
    box(s, x, 2.43, 3.93, 2.75, fill=LIGHT,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    for j, l in enumerate(lines):
        box(s, x + 0.35, 2.85 + j * 0.75, 0.15, 0.15, fill=GOLD)
        text(s, x + 0.66, 2.78 + j * 0.75, 3.05, 0.7,
             [{'runs': [(l, 15.5, False, INK)], 'line_spacing': 1.15}])
note_band(s, 5.5, '编号规则：数组下标即编号，parent < child，'
                  '推荐先序布局；允许共享子树。', h=0.85, size=15.5)
footer(s, 11)

# ================================================================
# Slide 12 · Trust boundary
# ================================================================
s = new_slide()
header(s, '设计核心：可信边界', 'Trust Boundary')
box(s, 0.62, 1.75, 5.15, 3.9, fill=RED_L,
    shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
text(s, 0.95, 2.05, 4.5, 0.55,
     [{'runs': [('不可信：搜索器输出', 19, True, RED)]}])
for i, t in enumerate(['C++ 搜索器（DFPN / VCF）', '启发式与剪枝',
                       '置换表与各种缓存', 'C++ 预检（仅诊断）']):
    box(s, 0.98, 2.8 + i * 0.68, 0.16, 0.16, fill=RED)
    text(s, 1.3, 2.73 + i * 0.68, 4.3, 0.62,
         [{'runs': [(t, 15.5, False, INK)]}])
box(s, 7.56, 1.75, 5.15, 3.9, fill=GREEN_L,
    shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
text(s, 7.89, 2.05, 4.5, 0.55,
     [{'runs': [('可信：Lean 检查器', 19, True, GREEN)]}])
for i, t in enumerate(['checkCertificate', '重算终局胜负',
                       '验证合法性与子局面', '对手全应手覆盖']):
    box(s, 7.92, 2.8 + i * 0.68, 0.16, 0.16, fill=GREEN)
    text(s, 8.24, 2.73 + i * 0.68, 4.3, 0.62,
         [{'runs': [(t, 15.5, False, INK)]}])
arrow(s, 5.77, 3.7, 7.56, 3.7, color=GRAY, width=2.6)
text(s, 5.68, 2.9, 2.0, 0.6,
     [{'runs': [('重新检查', 14, True, GRAY)], 'align': PP_ALIGN.CENTER}])
text(s, 0.95, 5.35, 4.5, 0.4,
     [{'runs': [('→ 只生成候选，可以出错', 14, True, RED)]}])
text(s, 7.89, 5.35, 4.6, 0.4,
     [{'runs': [('→ 通过 → CanForceWin 定理', 14, True, GREEN)]}])
note_band(s, 6.05, '搜索器可以犯错，检查器不放错；通过检查的证书才是定理。')
footer(s, 12)

# ================================================================
# Slide 13 · Results: how far did we get
# ================================================================
s = new_slide()
header(s, '已完成：做到什么程度', 'Results')
steps = [
    ('① 一步胜', True, ['四连局面', '证书检查通过']),
    ('② 两个应手', True, ['AND 根节点', '证书检查通过']),
    ('③ 有限深度', False, ['参考实现存在', '规模仍然有限']),
    ('④ 导出复检', True, ['C++ 导出 Lean 源码', '局部局面复检通过']),
]
x = 0.62
for i, (t, done, lines) in enumerate(steps):
    c = GREEN if done else ORANGE
    mark = '✓ 已完成' if done else '◐ 部分完成'
    sp = box(s, x, 1.7, 2.72, 0.72, fill=c,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
    text_in(sp, [{'runs': [(t, 17, True, WHITE)]}])
    box(s, x, 2.42, 2.72, 3.15, fill=LIGHT,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    box(s, x + 0.25, 2.62, 2.22, 0.5, fill=(GREEN_L if done else GOLD_L),
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.25)
    text_in(box(s, x + 0.25, 2.62, 2.22, 0.5,
                fill=(GREEN_L if done else GOLD_L),
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.25),
            [{'runs': [(mark, 13, True, (GREEN if done else GOLD_D))]}])
    for j, l in enumerate(lines):
        box(s, x + 0.25, 3.4 + j * 0.85, 0.15, 0.15, fill=GOLD)
        text(s, x + 0.55, 3.32 + j * 0.85, 2.05, 0.8,
             [{'runs': [(l, 14.5, False, INK)], 'line_spacing': 1.15}])
    if i < 3:
        arrow(s, x + 2.72, 2.06, x + 3.18, 2.06, color=GRAY, width=2.2)
    x += 3.18
note_band(s, 5.9, '全部成果以 Lean 复检通过为准：一步胜 / 双应手 / 开四 VCF。')
footer(s, 13)

# ================================================================
# Slide 14 · Chapter: verified draw positions
# ================================================================
s = new_slide(NAVY)
box(s, 0, 7.28, 13.333, 0.22, fill=GOLD)
text(s, 0.9, 1.6, 11.5, 0.5,
     [{'runs': [('04 · VERIFIED DRAW POSITIONS', 16, True, GOLD)]}])
text(s, 0.9, 2.3, 11.5, 1.5,
     [{'runs': [('9 个机器验证的 7×7 和棋局面', 44, True, WHITE)]}])
text(s, 0.9, 3.9, 11.0, 0.7,
     [{'runs': [('每个局面：21 黑 + 21 白 · 轮到黑方 · 7 个空格 · 双方都无法成五',
                 18, False, RGBColor(0xCF, 0xD8, 0xE8))]}])
box(s, 0.95, 4.9, 2.6, 0.05, fill=GOLD)
text(s, 0.9, 5.15, 11.5, 0.6,
     [{'runs': [('不可信搜索器找证书 → Lean 检查器逐节点重算 → StandardDraw 定理',
                 17, False, RGBColor(0xF2, 0xCE, 0x94))]}])
footer(s, 14)

# ================================================================
# Slide 15 · How a draw is verified
# ================================================================
s = new_slide()
header(s, '和棋是怎么验证的', 'Trust Chain')
steps = [
    ('① 构造局面', '每个长度 5 窗口都含黑白两子\n→ 任何一方永远无法成五', NAVY2),
    ('② C++ 找证书', '白防黑 + 黑防白 两张防守证书\n(搜索器不可信, 只出候选)', BLUE),
    ('③ Lean 重算', 'checkDefenseCertificateAt\n逐节点重算局面/轮次/终局/应手', GOLD_D),
    ('④ 组合出定理', '两张证书 + soundness\n⇒ StandardDraw 局面', GREEN),
]
x = 0.62
for i, (t1, t2, c) in enumerate(steps):
    sp = box(s, x, 1.7, 2.85, 1.55, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.1)
    text_in(sp, [{'runs': [(t1, 17, True, WHITE)]},
                 {'runs': [(t2, 11.5, False, WHITE)]}])
    if i < 3:
        conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x + 2.85), Inches(2.47),
                                      Inches(x + 3.28), Inches(2.47))
        conn.line.color.rgb = GRAY
        conn.line.width = Pt(2.0)
    x += 3.28
text(s, 0.62, 3.75, 12.1, 0.55,
     [{'runs': [('每个局面独立验证 4 项断言 (全部 native_decide, 不依赖 C++ 结果):',
                 16, True, NAVY)]}])
checks = [
    ('白防黑证书 checkDefenseCertificateAt = true', '9 / 9 通过'),
    ('黑防白证书 checkDefenseCertificateAt = true', '9 / 9 通过'),
    ('局面形状: 21黑 + 21白 + 7空 + 轮到黑方', '9 / 9 通过'),
    ('60 个长度 5 窗口全部双色', '9 / 9 通过'),
    ('StandardDraw <position> 由两张证书组合', '9 / 9 通过'),
]
y = 4.4
for i, (t, r) in enumerate(checks):
    box(s, 0.62, y, 8.6, 0.42, fill=LIGHT,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    text(s, 0.95, y, 8.2, 0.42, [{'runs': [(t, 14, False, INK)]}],
         anchor=MSO_ANCHOR.MIDDLE)
    chip = box(s, 9.5, y, 3.0, 0.42, fill=GREEN_L,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    text_in(chip, [{'runs': [(r, 13.5, True, GREEN)]}])
    y += 0.5
note = box(s, 0.62, 6.5, 12.09, 0.5, fill=GOLD_L,
           shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.2)
text_in(note, [{'runs': [('搜索器可以犯错, 检查器不放错; 通过检查的证书才是定理。',
                          14, True, GOLD_D)]}])
footer(s, 15)

# ================================================================
# Slides 16..24 · one per draw position (with board image)
# ================================================================
for idx, (lean, seed, png) in enumerate(DRAW_POSITIONS):
    page = 16 + idx
    s = new_slide()
    header(s, f'和棋局面 {idx + 1} / 9', f'{lean} · {seed}')
    img_path = os.path.join(BOARD_DIR, png)
    if os.path.exists(img_path):
        s.shapes.add_picture(img_path, Inches(0.75), Inches(1.35),
                             height=Inches(5.35))
    fx = 6.45
    facts = [
        ('局面构成', '21 黑 + 21 白 · 7 空格 · 轮到黑方'),
        ('结构性保证', '每个长度 5 窗口同时含黑白两子\n任何一方都无法成五'),
        ('白防黑证书', 'WhiteCanPreventBlackWin  √ 检查器通过'),
        ('黑防白证书', 'BlackCanPreventWhiteWin  √ 检查器通过'),
        ('和棋定理', 'StandardDraw <root>  √ 组合成立'),
    ]
    y = 1.5
    for i, (t1, t2) in enumerate(facts):
        box(s, fx, y, 6.1, 1.02, fill=LIGHT,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
        text(s, fx + 0.3, y + 0.12, 2.3, 0.4,
             [{'runs': [(t1, 15.5, True, BLUE)]}])
        text(s, fx + 0.3, y + 0.48, 5.6, 0.5,
             [{'runs': [(t2, 13, False, INK)], 'line_spacing': 1.1}])
        y += 1.13
    footer(s, page)

# ================================================================
# Slide 25 · Draw positions summary table
# ================================================================
s = new_slide()
header(s, '全部 9 个局面: 汇总', 'Summary')
col_x = [0.62, 3.1, 5.2, 7.0, 8.9, 10.9]
col_w = [2.48, 2.1, 1.8, 1.9, 2.0, 1.6]
heads = ['Lean 定理文件', '种子', '白防证书', '黑防证书', 'StandardDraw', '窗口双色']
for i, h in enumerate(heads):
    box(s, col_x[i], 1.5, col_w[i], 0.55, fill=NAVY)
    text(s, col_x[i], 1.5, col_w[i], 0.55,
         [{'runs': [(h, 14, True, WHITE)], 'align': PP_ALIGN.CENTER}],
         anchor=MSO_ANCHOR.MIDDLE)
y = 2.15
for idx, (lean, seed, png) in enumerate(DRAW_POSITIONS):
    fill = LIGHT if idx % 2 == 0 else WHITE
    for i, val in enumerate([lean, seed, '✓', '✓', '✓ 定理', '✓']):
        box(s, col_x[i], y, col_w[i], 0.52, fill=fill)
        color = GREEN if val == '✓' or val == '✓ 定理' else INK
        text(s, col_x[i], y, col_w[i], 0.52,
             [{'runs': [(val, 13, (val != lean), color)],
               'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.52
band = box(s, 0.62, 6.35, 12.09, 0.6, fill=GREEN_L,
           shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
text_in(band, [{'runs': [('9 / 9 全部通过: 每张证书经 Lean 检查器重算, 每个局面都证明了 '
                          'StandardDraw。', 15, True, GREEN)]}])
footer(s, 25)

# ================================================================
# Slide 26 · Gaps 1: functional
# ================================================================
s = new_slide()
header(s, '不足 ①：功能缺口', 'Not Done Well · Function')
issues = [
    ('空棋盘开局求解 · 未完成',
     '没有从 7×7 空棋盘开始的必胜/和棋证书，不能宣称 initial_black_wins'),
    ('空棋盘和棋 · 未完成',
     '中盘 9 局面的和棋已验证；空棋盘 StandardDraw 仍是开放目标'),
    ('缓存 · 未接入',
     'SearchMemo / EngineMemo 已定义，但还没有接入引擎默认搜索流程'),
    ('性能 · 未验证',
     '目前只有一步胜、双应手、开四 VCF 与中盘和棋等示例，规模没有跑起来'),
]
for i, (t, d) in enumerate(issues):
    x = 0.62 + (i % 2) * 6.2
    y = 1.7 + (i // 2) * 2.35
    issue_card(s, x, y, 5.9, 2.1, t, d)
note_band(s, 6.35, '这些缺口意味着：项目目前给不出“空棋盘黑棋先手必胜”或“空棋盘和棋”的结论。',
          fill=RED_L, color=RED)
footer(s, 26)

# ================================================================
# Slide 27 · Gaps 2: engineering
# ================================================================
s = new_slide()
header(s, '不足 ②：工程与验证', 'Not Done Well · Engineering')
issues = [
    ('参数没有系统调优',
     'maxDepth / maxVcfDepth / maxNodes 用的是保守默认值，没做过参数扫描'),
    ('缺少基准数据',
     '证书大小、Lean 检查耗时、缓存命中率都没有系统记录和对比'),
    ('导出→复检链路粗糙',
     'C++ 导出 Lean 源码后要手工跑检查，自动化程度不高'),
    ('投入偏重正确性',
     '时间主要花在规则与检查器上，搜索规模与性能投入不足'),
]
for i, (t, d) in enumerate(issues):
    x = 0.62 + (i % 2) * 6.2
    y = 1.7 + (i // 2) * 2.35
    issue_card(s, x, y, 5.9, 2.1, t, d)
note_band(s, 6.35, '一句话：框架完整、正确性优先，但离完整 7×7 求解还有明显距离。',
          fill=RED_L, color=RED)
footer(s, 27)

# ================================================================
# Slide 28 · Closing
# ================================================================
s = new_slide(NAVY)
box(s, 0, 7.28, 13.333, 0.22, fill=GOLD)
text(s, 0.9, 1.8, 11.5, 1.3,
     [{'runs': [('谢谢', 60, True, WHITE)], 'align': PP_ALIGN.CENTER}])
box(s, 6.17, 3.35, 1.0, 0.05, fill=GOLD)
text(s, 0.9, 3.7, 11.5, 0.7,
     [{'runs': [('搜索器负责“找”，Lean 负责“证”', 24, False,
                 RGBColor(0xF2, 0xCE, 0x94))], 'align': PP_ALIGN.CENTER}])
text(s, 0.9, 4.7, 11.5, 0.55,
     [{'runs': [('已做：规则 · 证书检查链 · 搜索器与示例复检 · 9 个中盘和棋局面', 16, False,
                 RGBColor(0xCF, 0xD8, 0xE8))], 'align': PP_ALIGN.CENTER}])
text(s, 0.9, 5.3, 11.5, 0.55,
     [{'runs': [('未做：空棋盘开局求解 · 空棋盘和棋 · 大规模验证', 16, False,
                 RGBColor(0xCF, 0xD8, 0xE8))], 'align': PP_ALIGN.CENTER}])

prs.core_properties.title = '五子棋形式化项目阶段汇报'
prs.core_properties.author = 'Gomoku Formalization'

# ---------------- speaker notes (one per slide, in slide order) ----------------
NOTES = [
    # 1 封面
    '各位好，我汇报五子棋形式化项目的阶段进展。核心一句话：让一个不可信的 C++ 搜索器'
    '去找必胜策略的候选，再由 Lean 检查器重新验证；只有通过检查的，才是定理。',
    # 2 目录
    '汇报分四部分：先讲项目是什么，再讲搜索器和验证怎么设计，然后讲目前做到什么程度，'
    '最后重点讲哪些地方还做得不好。',
    # 3 项目是什么
    '项目用 Lean 4 形式化固定规则的五子棋：7×7 棋盘、五连即胜、黑棋先手、无禁手，'
    '长连也算胜。我们做三件事：形式化规则本身、写搜索器找必胜策略、让所有结论都经过'
    'Lean 检查器确认。',
    # 4 为什么需要搜索器
    '为什么需要搜索器？49 个点的博弈树是天文数字，手工推演不可能。但程序又会犯错，'
    '不能把程序结果直接当结论。于是分工：搜索器找候选，Lean 出定理。',
    # 5 搜索语义
    '搜索器要找的是一棵必胜策略树。目标方回合是 OR 节点，找到一个能赢的着法就行；'
    '对手回合是 AND 节点，每个合法应手都要有胜子树，漏一个就失败。叶子必须是真实胜局，'
    '和棋、输棋、搜索耗尽都不算赢。',
    # 6 整体架构
    '搜索器用 C++ 实现：输入局面后，DFPN 主搜索负责证明数搜索，VCF 快攻负责连续冲四'
    '检测，产出候选证书，再导出成 Lean 源码。中间还有置换表、迭代加深、资源上限这些'
    '组件。注意：整条流水线产出的都只是候选数据。',
    # 7 DFPN
    'DFPN 是核心算法。每个节点记录两个数：证明数 pn，证明它必胜还要展开多少节点；'
    '反证数 dn，证明它必败还要多少节点。每次优先展开 pn 或 dn 最小的节点，把预算花在'
    '最可能的分支上。它深度优先、内存友好，还能直接输出证明路径。',
    # 8 VCF
    'VCF 是连续冲四快攻。冲四就是再下一子就成五，对手必须立刻防守；连续冲四直到成五，'
    '就是一条杀棋路线。VCF 有独立预算，攻不下来就退回 DFPN 做全局面搜索。注意：'
    'VCF 找到的路线依然是候选，要过检查器。',
    # 9 加速技巧
    '两个加速技巧。威胁排序：着法按威胁从高到低排，成五优先、堵五其次，先搜最有希望'
    '的分支；Lean 证明过排序不会引入非法着法。强制防守剪枝：目标方回合如果已经有成五'
    '或堵五的着法，就只搜这些。两者只影响速度，不影响正确性。',
    # 10 资源控制
    '搜索有各种资源上限：深度、节点数、置换表大小、证书节点数。状态分五种：found 是'
    '找到候选证书，其余都是没找到。最关键的纪律：搜不到不等于必败，只能说本次资源下'
    '没有找到证据，绝不能写成和棋或必败。',
    # 11 证书格式
    '搜索器的输出是 CompactCertificate，三种节点：terminal 是胜利叶子，proverMove 是'
    '目标方选一个着法，opponentMoves 是列出对手全部应手。编号规则是父节点必须小于'
    '子节点，自环和回边一律拒绝。',
    # 12 可信边界
    '这是整个设计最核心的地方。左边这些——搜索算法、启发式、剪枝、缓存、C++ 预检——'
    '全部不可信，只生成候选。右边 Lean 检查器把每个节点重新检查：重算终局、验证合法'
    '着法、验证子局面、检查对手全覆盖，通过之后才得到 CanForceWin 定理。',
    # 13 已完成
    '验收按四步走。一步胜、两个应手、导出复检三步已经打通，证书都在 Lean 里复检通过；'
    '有限深度搜索有参考实现，但规模还比较有限。整体上框架是完整的。',
    # 14 和棋章节封面
    '接下来进入第四个部分：9 个机器验证的 7×7 和棋局面。核心纪律不变：'
    'C++ 搜索器只负责找候选证书，是否成立完全由 Lean 检查器决定。',
    # 15 和棋怎么验证
    '每个局面的验证分四步：先构造一个每个长度5窗口都同时含黑白两子的局面，'
    '然后 C++ 分别搜索白防黑和黑防白两张防守证书，Lean 检查器逐节点重算，'
    '最后用 soundness 定理组合出 StandardDraw。九个局面、五项断言全部通过。',
    # 16 局面1
    '局面一：Draw7x7，最初的种子。21黑21白，7个空格，轮到黑方。',
    # 17 局面2
    '局面二：Draw7x7s2，随机化移除顺序得到的第二个种子。',
    # 18 局面3
    '局面三：Draw7x7s3。',
    # 19 局面4
    '局面四：Draw7x7s4。',
    # 20 局面5
    '局面五：Draw7x7s5。',
    # 21 局面6
    '局面六：Draw7x7s6。',
    # 22 局面7
    '局面七：Draw7x7s7。',
    # 23 局面8
    '局面八：Draw7x7s8。',
    # 24 局面9
    '局面九：Draw7x7s9。',
    # 25 汇总
    '汇总：九个局面的两张证书全部通过 Lean 检查，每个局面都得到了 StandardDraw 定理。'
    '需要说明：这些是特定的中盘局面，空棋盘的和棋定理仍然是开放目标。',
    # 26 不足①
    '但必须坦率地说，还有明显缺口。第一，空棋盘开局求解没有完成，不能宣称黑棋先手必胜；'
    '第二，空棋盘和棋还没有证明，目前只有中盘局面的和棋；第三，缓存已经定义好，但还没'
    '接入默认搜索流程；第四，目前只在一步胜、双应手和中盘和棋这些局面上验证过，'
    '规模没有跑起来。',
    # 27 不足②
    '工程上也有问题：搜索参数用的是保守默认值，没有系统调优；证书大小、检查耗时、'
    '缓存命中率这些基准数据都没有记录；导出到复检的链路还要手工跑，自动化不够；'
    '时间主要花在保证正确性上，性能投入不足。一句话：框架完整，但离完整 7×7 求解'
    '还有明显距离。',
    # 28 结束
    '总结：搜索器负责找，Lean 负责证。已经做的是规则、证书检查链、搜索器、示例复检，'
    '以及 9 个中盘和棋局面；没有做的是空棋盘开局求解、空棋盘和棋、大规模验证。'
    '谢谢大家，欢迎提问。',
]
for slide, note in zip(prs.slides, NOTES):
    slide.notes_slide.notes_text_frame.text = note

OUT = r'C:\Users\lenovo\Desktop\GomokuFormalization\五子棋项目阶段汇报.pptx'
try:
    prs.save(OUT)
    print('saved:', OUT, 'slides =', len(prs.slides._sldIdLst))
except PermissionError:
    OUT = r'C:\Users\lenovo\Desktop\GomokuFormalization\五子棋项目阶段汇报_新版.pptx'
    prs.save(OUT)
    print('original file locked; saved as:', OUT,
          'slides =', len(prs.slides._sldIdLst))
